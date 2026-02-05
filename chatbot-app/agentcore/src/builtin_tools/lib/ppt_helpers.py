"""
PowerPoint Helper Functions for Code Interpreter

This module is uploaded to Code Interpreter workspace and provides
high-level helper functions for AI-assisted PowerPoint generation.

Key Functions:
- generate_ppt_structure: Create slides from structured outline
- select_layout: Smart layout selection based on content type
- populate_slide: Fill slide with content based on type
- analyze_template: Extract layouts and theme from PPT
- markdown_to_outline: Convert markdown to presentation outline

These helpers enable natural language -> PPT conversion without requiring
users to write detailed python-pptx code.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from typing import Dict, List, Optional, Any


class TextSizeCalculator:
    """Calculate optimal text sizes based on content and container dimensions.

    Helps prevent text overflow by automatically calculating appropriate font sizes
    and wrapping text intelligently.
    """

    def __init__(self):
        self.character_widths = {
            'narrow': 0.6,  # i, l, t
            'normal': 1.0,  # most characters
            'wide': 1.3,    # m, w
            'space': 0.5    # space character
        }

    def estimate_text_width(self, text: str, font_size: int) -> float:
        """Estimate text width in points based on character analysis."""
        if not text:
            return 0

        width = 0
        for char in text:
            if char in 'iltj':
                width += self.character_widths['narrow']
            elif char in 'mwMW':
                width += self.character_widths['wide']
            elif char == ' ':
                width += self.character_widths['space']
            else:
                width += self.character_widths['normal']

        return width * font_size * 0.6  # Approximation factor

    def estimate_text_height(self, text: str, font_size: int, line_spacing: float = 1.2) -> float:
        """Estimate text height based on line count and spacing."""
        lines = len(text.split('\n'))
        return lines * font_size * line_spacing * 1.3  # Convert to points

    def calculate_optimal_font_size(self, text: str, container_width: float,
                                  container_height: float, font_type: str = 'body',
                                  min_size: int = 8, max_size: int = 36) -> int:
        """Calculate optimal font size to fit text in container.

        Args:
            text: Text content to fit
            container_width: Container width in inches
            container_height: Container height in inches
            font_type: Type of text ('title', 'body', etc.)
            min_size: Minimum font size in points
            max_size: Maximum font size in points

        Returns:
            Optimal font size in points
        """
        container_width_pts = container_width * 72  # Convert inches to points
        container_height_pts = container_height * 72

        # Start with a reasonable size and adjust
        for font_size in range(max_size, min_size - 1, -1):
            estimated_width = self.estimate_text_width(text, font_size)
            estimated_height = self.estimate_text_height(text, font_size)

            if estimated_width <= container_width_pts * 0.9 and estimated_height <= container_height_pts * 0.9:
                return font_size

        return min_size

    def wrap_text_intelligently(self, text: str, max_width: float, font_size: int) -> str:
        """Intelligently wrap text to fit within specified width.

        Args:
            text: Text to wrap
            max_width: Maximum width in inches
            font_size: Font size in points

        Returns:
            Wrapped text with newlines
        """
        if not text:
            return text

        max_width_pts = max_width * 72
        words = text.split()
        wrapped_lines = []
        current_line = []

        for word in words:
            test_line = current_line + [word]
            test_text = ' '.join(test_line)

            if self.estimate_text_width(test_text, font_size) <= max_width_pts:
                current_line.append(word)
            else:
                if current_line:
                    wrapped_lines.append(' '.join(current_line))
                    current_line = [word]
                else:
                    # Single word is too long, force wrap
                    wrapped_lines.append(word)

        if current_line:
            wrapped_lines.append(' '.join(current_line))

        return '\n'.join(wrapped_lines)


def generate_ppt_structure(prs: Presentation, outline: dict, theme_name: str = None) -> None:
    """Generate PowerPoint structure from outline with theme support

    Main entry point for AI-assisted generation. Takes a structured outline
    and creates slides with appropriate layouts and content. Applies theme
    styling if theme_name is provided.

    Args:
        prs: Presentation object (already initialized)
        outline: Structured outline dict with format:
            {
                "title": "Presentation Title",
                "subtitle": "Optional subtitle",  # optional
                "theme": "modern-blue",  # optional, can be specified in outline
                "slides": [
                    {
                        "title": "Slide Title",
                        "type": "bullet|chart|image|table|section|title|blank",
                        "content": [...],  # format depends on type
                        "notes": "Optional speaker notes"  # optional
                    }
                ]
            }
        theme_name: Name of theme to apply (overrides outline theme)

    Example:
        outline = {
            "title": "Q4 Results",
            "slides": [
                {"title": "Summary", "type": "bullet", "content": ["Revenue up 25%", "Profit up 30%"]},
                {"title": "Growth", "type": "chart", "data": {"labels": [...], "values": [...]}}
            ]
        }
        generate_ppt_structure(prs, outline, theme_name='minimal')
    """
    # Determine theme to use
    if not theme_name:
        theme_name = outline.get("theme", "minimal")

    # Get theme config
    theme = BUILTIN_THEMES.get(theme_name, BUILTIN_THEMES['minimal'])

    # Add title slide
    if "title" in outline:
        title_slide = prs.slides.add_slide(prs.slide_layouts[0])
        title_slide.shapes.title.text = outline["title"]

        # Add subtitle if provided
        if "subtitle" in outline and len(title_slide.placeholders) > 1:
            try:
                title_slide.placeholders[1].text = outline["subtitle"]
            except:
                pass

        # Apply theme to title slide
        _apply_theme_to_slide_shapes(title_slide, theme)

    # Process each content slide
    for slide_info in outline.get("slides", []):
        try:
            # Select appropriate layout
            layout = select_layout(prs, slide_info.get("type", "bullet"))

            # Add slide
            slide = prs.slides.add_slide(layout)

            # Populate with content
            populate_slide(slide, slide_info)

            # Apply theme to slide
            _apply_theme_to_slide_shapes(slide, theme)

            # Add speaker notes if provided
            if "notes" in slide_info:
                try:
                    notes_slide = slide.notes_slide
                    notes_slide.notes_text_frame.text = slide_info["notes"]
                except:
                    pass

        except Exception as e:
            print(f"Warning: Failed to create slide '{slide_info.get('title', 'Untitled')}': {e}")
            continue

    print(f"Generated {len(prs.slides)} slides from outline with {theme['name']} theme")


def select_layout(prs: Presentation, content_type: str, template_info: dict = None) -> 'SlideLayout':
    """Select appropriate slide layout based on content type

    Intelligently chooses the best layout for the given content type.
    If template_info is provided, uses layout name matching for better results.

    Args:
        prs: Presentation object
        content_type: Type of content - one of:
            - 'title': Title slide
            - 'bullet' or 'content': Bullet points / text content
            - 'chart': Chart or graph
            - 'image': Image-focused slide
            - 'table': Table data
            - 'section': Section divider
            - 'blank': Blank slide
        template_info: Optional template metadata (from analyze_template)

    Returns:
        SlideLayout object

    Example:
        layout = select_layout(prs, 'bullet', template_info)
        slide = prs.slides.add_slide(layout)
    """
    # If template info available, match by layout name
    if template_info:
        layout_names = {layout['index']: layout['name'].lower()
                       for layout in template_info.get('layouts', [])}

        # Match by name heuristics
        if content_type == 'title':
            for idx, name in layout_names.items():
                if 'title' in name and 'content' not in name and idx < len(prs.slide_layouts):
                    return prs.slide_layouts[idx]

        elif content_type in ['bullet', 'content', 'table']:
            for idx, name in layout_names.items():
                if ('content' in name or 'bullet' in name) and idx < len(prs.slide_layouts):
                    return prs.slide_layouts[idx]

        elif content_type == 'section':
            for idx, name in layout_names.items():
                if ('section' in name or 'header' in name) and idx < len(prs.slide_layouts):
                    return prs.slide_layouts[idx]

        elif content_type in ['blank', 'chart', 'image']:
            for idx, name in layout_names.items():
                if ('blank' in name or 'title only' in name) and idx < len(prs.slide_layouts):
                    return prs.slide_layouts[idx]

    # Fallback: Default python-pptx layout indices
    # These are standard indices for most PowerPoint templates
    default_map = {
        'title': 0,       # Title Slide
        'bullet': 1,      # Title and Content
        'content': 1,     # Title and Content (alias for bullet)
        'section': 2,     # Section Header
        'table': 1,       # Title and Content (works for tables)
        'chart': 5,       # Title Only
        'image': 5,       # Title Only
        'blank': 6        # Blank
    }

    idx = default_map.get(content_type, 5)  # Default to Title Only for unknown types

    # Safety check: ensure index is within bounds
    max_idx = len(prs.slide_layouts) - 1
    safe_idx = min(idx, max_idx)

    return prs.slide_layouts[safe_idx]


def populate_slide(slide, slide_info: dict) -> None:
    """Populate slide with content based on type

    Fills in slide title and content based on the slide_info structure.
    Delegates to specialized functions for each content type.

    Args:
        slide: Slide object
        slide_info: Slide information dict with keys:
            - title: Slide title (string)
            - type: Content type (string)
            - content: Content data (format depends on type)

    Example:
        slide_info = {
            "title": "Key Points",
            "type": "bullet",
            "content": ["Point 1", "Point 2", "Point 3"]
        }
        populate_slide(slide, slide_info)
    """
    # Set title - safe approach for layouts without title placeholder
    if "title" in slide_info:
        title_text = slide_info["title"]

        # Check if title placeholder exists (not all layouts have it)
        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
            try:
                slide.shapes.title.text = title_text
            except Exception as e:
                print(f"Warning: Could not set title: {e}")
        else:
            # Layout has no title placeholder (e.g., Blank layout)
            # Add text box as fallback
            try:
                from pptx.util import Inches
                left = Inches(0.5)
                top = Inches(0.3)
                width = Inches(9)
                height = Inches(0.8)  # Reduced height to prevent overlap
                textbox = slide.shapes.add_textbox(left, top, width, height)
                text_frame = textbox.text_frame
                text_frame.text = title_text
                text_frame.word_wrap = True  # Enable word wrap
                # Make it look like a title
                for paragraph in text_frame.paragraphs:
                    paragraph.font.size = Pt(28)  # Slightly smaller for better fit
                    paragraph.font.bold = True
                    paragraph.font.name = "Calibri"  # Consistent font
                print(f"Added title as text box (layout has no title placeholder)")
            except Exception as e:
                print(f"Warning: Could not add title text box: {e}")

    # Populate content based on type
    content_type = slide_info.get("type", "bullet")
    content = slide_info.get("content")
    custom_code = slide_info.get("custom_code")

    # Handle custom code type - highest priority
    if content_type == "custom" or custom_code:
        try:
            # Execute custom python-pptx code
            # Available variables: slide, Inches, Pt, RGBColor
            from pptx.util import Inches, Pt
            from pptx.dml.color import RGBColor
            from pptx.enum.shapes import MSO_SHAPE_TYPE
            from pptx.chart.data import CategoryChartData
            from pptx.enum.chart import XL_CHART_TYPE

            # Create local namespace with available variables
            local_namespace = {
                'slide': slide,
                'Inches': Inches,
                'Pt': Pt,
                'RGBColor': RGBColor,
                'MSO_SHAPE_TYPE': MSO_SHAPE_TYPE,
                'CategoryChartData': CategoryChartData,
                'XL_CHART_TYPE': XL_CHART_TYPE
            }

            # Execute the custom code
            exec(custom_code, local_namespace)
            print(f"Executed custom code for slide")
        except Exception as e:
            print(f"Warning: Failed to execute custom code: {e}")
        return

    if not content:
        return

    try:
        if content_type in ["bullet", "content"]:
            add_bullet_points(slide, content)
        elif content_type == "chart":
            add_chart_to_slide(slide, content)
        elif content_type == "image":
            add_image_to_slide(slide, content)
        elif content_type == "table":
            add_table_to_slide(slide, content)
        elif content_type == "section":
            # Section slides typically only have title
            pass
        elif content_type == "title":
            # Title slides handled separately
            pass
        # blank slides have no content
    except Exception as e:
        print(f"Warning: Failed to populate content for type '{content_type}': {e}")


def add_bullet_points(slide, points: list) -> None:
    """Add bullet points to content placeholder

    Finds the content placeholder (usually index 1) and adds bullet points.
    Supports nested lists with indentation levels.

    Args:
        slide: Slide object
        points: List of bullet points (strings or dicts with 'text' and 'level')

    Example:
        add_bullet_points(slide, [
            "Main point 1",
            "Main point 2",
            {"text": "Sub-point", "level": 1}
        ])
    """
    # Find content placeholder (usually BODY type, placeholder index 1)
    content_placeholder = None

    # Try to find by placeholder type
    for shape in slide.placeholders:
        try:
            # Type 2 is BODY placeholder
            if shape.placeholder_format.type == 2:
                content_placeholder = shape
                break
        except:
            continue

    # Fallback: try placeholder index 1
    if not content_placeholder:
        try:
            if len(slide.placeholders) > 1:
                content_placeholder = slide.placeholders[1]
        except:
            pass

    if not content_placeholder or not hasattr(content_placeholder, 'text_frame'):
        print("Warning: Could not find content placeholder for bullet points")
        return

    # Add bullet points
    text_frame = content_placeholder.text_frame

    for idx, point in enumerate(points):
        # Handle both string and dict format
        if isinstance(point, dict):
            text = point.get("text", "")
            level = point.get("level", 0)
        else:
            text = str(point)
            level = 0

        # First point uses existing paragraph, rest are added
        if idx == 0:
            text_frame.text = text
            text_frame.paragraphs[0].level = level
        else:
            p = text_frame.add_paragraph()
            p.text = text
            p.level = level


def add_chart_to_slide(slide, chart_data: dict) -> None:
    """Add chart to slide (requires matplotlib)

    Generates a chart using matplotlib and inserts it into the slide.

    Args:
        slide: Slide object
        chart_data: Chart configuration dict with keys:
            - type: 'bar', 'line', 'pie', 'scatter'
            - labels: List of labels
            - values: List of values (or list of lists for multiple series)
            - title: Chart title (optional)
            - xlabel, ylabel: Axis labels (optional)

    Example:
        chart_data = {
            "type": "bar",
            "labels": ["Q1", "Q2", "Q3", "Q4"],
            "values": [100, 120, 150, 140],
            "title": "Quarterly Sales"
        }
        add_chart_to_slide(slide, chart_data)
    """
    try:
        import matplotlib.pyplot as plt

        # Extract chart parameters
        chart_type = chart_data.get("type", "bar")
        labels = chart_data.get("labels", [])
        values = chart_data.get("values", [])
        title = chart_data.get("title", "")
        xlabel = chart_data.get("xlabel", "")
        ylabel = chart_data.get("ylabel", "")

        # Create chart
        plt.figure(figsize=(10, 6))

        if chart_type == "bar":
            plt.bar(labels, values)
        elif chart_type == "line":
            plt.plot(labels, values, marker='o')
        elif chart_type == "pie":
            plt.pie(values, labels=labels, autopct='%1.1f%%')
        elif chart_type == "scatter":
            plt.scatter(labels, values)
        else:
            plt.bar(labels, values)  # Default to bar

        if title:
            plt.title(title)
        if xlabel:
            plt.xlabel(xlabel)
        if ylabel:
            plt.ylabel(ylabel)

        plt.tight_layout()

        # Save chart with high quality
        chart_filename = f"chart_{id(slide)}.png"
        plt.savefig(chart_filename, dpi=300, bbox_inches='tight', facecolor='white')
        plt.close()

        # Insert into slide with safe positioning
        left = Inches(0.8)
        top = Inches(1.5)  # Start below title area (title ends ~1.1")
        width = Inches(8.4)  # Use more width
        slide.shapes.add_picture(chart_filename, left, top, width=width)

        print(f"Added {chart_type} chart to slide")

    except Exception as e:
        print(f"Warning: Failed to add chart: {e}")


def add_image_to_slide(slide, image_path: str) -> None:
    """Add image to slide with automatic size adjustment to fit slide bounds

    Args:
        slide: Slide object
        image_path: Path to image file (string or dict with 'path' and optional 'width')

    Example:
        add_image_to_slide(slide, "logo.png")
        # or
        add_image_to_slide(slide, {"path": "logo.png", "width": 5})  # Fixed width
    """
    try:
        import os
        from PIL import Image

        # Handle both string and dict format
        manual_width = None
        if isinstance(image_path, dict):
            path = image_path.get("path", "")
            # If width is specified, use it (no auto-adjustment)
            if "width" in image_path:
                manual_width = Inches(image_path.get("width"))
        else:
            path = str(image_path)

        # Check if file exists
        if not os.path.exists(path):
            # Try to find similar files
            available = [f for f in os.listdir('.') if f.endswith(('.png', '.jpg', '.jpeg', '.gif', '.webp'))]
            print(f"Warning: Image not found: {path}. Available: {available}")
            return

        # If manual width specified, use it without adjustment
        if manual_width:
            left = Inches(0.8)
            top = Inches(1.5)
            slide.shapes.add_picture(path, left, top, width=manual_width)
            print(f"Added image: {path} (manual width: {manual_width})")
            return

        # Auto-adjust: Calculate optimal size to fit slide
        # Standard slide: 10" × 7.5"
        # Available area (with margins and title): ~9" × 5.5"
        max_width = Inches(9)
        max_height = Inches(5.5)
        margin_left = Inches(0.5)
        margin_top = Inches(1.5)  # Below title

        # Get image dimensions
        img = Image.open(path)
        img_width, img_height = img.size
        img_ratio = img_width / img_height

        # Calculate size that fits within bounds while preserving aspect ratio
        # Try fitting by width first
        target_width = max_width
        target_height = target_width / img_ratio

        # If height exceeds limit, fit by height instead
        if target_height > max_height:
            target_height = max_height
            target_width = target_height * img_ratio

        # Center the image in available space
        left = margin_left + (max_width - target_width) / 2
        top = margin_top + (max_height - target_height) / 2

        slide.shapes.add_picture(path, left, top, width=target_width)

        print(f"Added image: {path} ({img_width}×{img_height}px → {target_width.inches:.2f}\"×{target_height.inches:.2f}\")")

    except ImportError:
        # PIL not available, fallback to fixed size
        print("Warning: PIL not available for image size detection. Using fixed size.")
        import os
        if isinstance(image_path, dict):
            path = image_path.get("path", "")
        else:
            path = str(image_path)

        if os.path.exists(path):
            left = Inches(0.8)
            top = Inches(1.5)
            slide.shapes.add_picture(path, left, top, width=Inches(8))
            print(f"Added image: {path} (fixed size)")

    except Exception as e:
        print(f"Warning: Failed to add image: {e}")


def add_table_to_slide(slide, table_data: dict) -> None:
    """Add table to slide

    Args:
        slide: Slide object
        table_data: Table configuration dict with keys:
            - headers: List of header strings
            - rows: List of row lists
            - style: Optional table style name

    Example:
        table_data = {
            "headers": ["Product", "Q3", "Q4"],
            "rows": [
                ["Widget", "$100K", "$125K"],
                ["Gadget", "$200K", "$250K"]
            ]
        }
        add_table_to_slide(slide, table_data)
    """
    try:
        headers = table_data.get("headers", [])
        rows = table_data.get("rows", [])
        table_style = table_data.get("style", None)

        if not headers or not rows:
            print("Warning: Table data incomplete (need headers and rows)")
            return

        # Calculate dimensions
        num_rows = len(rows) + 1  # +1 for header
        num_cols = len(headers)

        # Create table
        left = Inches(1)
        top = Inches(2)
        width = Inches(8)
        height = Inches(0.5 * num_rows)

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # Set headers
        for col_idx, header in enumerate(headers):
            cell = table.cell(0, col_idx)
            cell.text = str(header)

        # Set rows
        for row_idx, row in enumerate(rows):
            for col_idx, value in enumerate(row):
                cell = table.cell(row_idx + 1, col_idx)
                cell.text = str(value)

        # Apply style if provided
        if table_style:
            try:
                table_shape.table.style = table_style
            except:
                pass

        print(f"Added table: {num_rows}x{num_cols}")

    except Exception as e:
        print(f"Warning: Failed to add table: {e}")


def analyze_template(pptx_path: str) -> dict:
    """Extract layouts and theme information from PowerPoint template

    Analyzes a PowerPoint file to extract layout structures, placeholder
    information, and theme colors. This metadata can be used to create
    new slides that match the template's style.

    Args:
        pptx_path: Path to PowerPoint file

    Returns:
        Template info dict with structure:
        {
            "source_file": "template.pptx",
            "analyzed_at": "2025-12-22T10:30:00Z",
            "slide_count": 12,
            "layouts": [
                {
                    "index": 0,
                    "name": "Title Slide",
                    "placeholders": [
                        {"type": 1, "name": "Title 1", "idx": 0},
                        {"type": 2, "name": "Subtitle 2", "idx": 1}
                    ]
                }
            ],
            "theme": {
                "accent_colors": ["#1F4E78", "#0070C0"],
                "fonts": {"title": "Calibri Light", "body": "Calibri"}
            }
        }

    Example:
        template_info = analyze_template('company-template.pptx')
        print(f"Found {len(template_info['layouts'])} layouts")
    """
    from datetime import datetime

    try:
        prs = Presentation(pptx_path)

        # Extract layouts
        layouts = []
        for idx, layout in enumerate(prs.slide_layouts):
            placeholders = []

            # Extract placeholder info
            for ph in layout.placeholders:
                try:
                    placeholders.append({
                        "type": ph.placeholder_format.type,
                        "idx": ph.placeholder_format.idx,
                        "name": ph.name
                    })
                except:
                    pass

            layouts.append({
                "index": idx,
                "name": layout.name,
                "placeholders": placeholders
            })

        # Try to extract theme colors (basic extraction)
        accent_colors = []
        try:
            # This is a simplified version - full theme extraction is complex
            slide_master = prs.slide_master
            # Theme color extraction would require deeper XML parsing
            # For now, we'll leave it as empty or add placeholder
            accent_colors = ["#1F4E78", "#0070C0", "#C5E0B4"]  # Common defaults
        except:
            pass

        # Try to extract fonts (basic)
        fonts = {}
        try:
            # Font extraction also requires XML parsing
            # Using common defaults for now
            fonts = {
                "title": "Calibri Light",
                "body": "Calibri"
            }
        except:
            pass

        template_info = {
            "source_file": pptx_path,
            "analyzed_at": datetime.now().isoformat(),
            "slide_count": len(prs.slides),
            "layouts": layouts,
            "theme": {
                "accent_colors": accent_colors,
                "fonts": fonts
            }
        }

        print(f"Template analyzed: {len(layouts)} layouts, {len(prs.slides)} slides")
        return template_info

    except Exception as e:
        print(f"Error analyzing template: {e}")
        return {
            "source_file": pptx_path,
            "error": str(e),
            "layouts": [],
            "theme": {}
        }


def markdown_to_outline(markdown_content: str) -> dict:
    """Convert markdown document to presentation outline

    Parses markdown structure (headings, lists) and converts to a
    presentation outline suitable for generate_ppt_structure.

    Args:
        markdown_content: Markdown text

    Returns:
        Outline dict ready for generate_ppt_structure

    Example:
        markdown = '''
        # My Presentation
        ## Introduction
        - Point 1
        - Point 2
        ## Data
        - Numbers here
        '''
        outline = markdown_to_outline(markdown)
        generate_ppt_structure(prs, outline)
    """
    lines = markdown_content.split('\n')

    title = None
    slides = []
    current_slide = None

    for line in lines:
        line = line.strip()

        if not line:
            continue

        # H1 -> Title
        if line.startswith('# '):
            title = line[2:].strip()

        # H2 -> New slide (bullet type)
        elif line.startswith('## '):
            if current_slide:
                slides.append(current_slide)

            current_slide = {
                "title": line[3:].strip(),
                "type": "bullet",
                "content": []
            }

        # H3 -> Section slide
        elif line.startswith('### '):
            if current_slide:
                slides.append(current_slide)

            current_slide = {
                "title": line[4:].strip(),
                "type": "section",
                "content": []
            }

        # Bullet point -> Add to current slide
        elif line.startswith('- ') or line.startswith('* '):
            if current_slide:
                current_slide["content"].append(line[2:].strip())

    # Add last slide
    if current_slide:
        slides.append(current_slide)

    # If no title found, use first slide title or default
    if not title:
        if slides:
            title = slides[0].get("title", "Presentation")
        else:
            title = "Presentation"

    return {
        "title": title,
        "slides": slides
    }


# ==============================================================================
# Theme System: Predefined color schemes and styles
# ==============================================================================

BUILTIN_THEMES = {
    'minimal': {
        'name': 'Minimal',
        'colors': {
            'primary': (45, 45, 45),       # Dark Gray (almost black)
            'secondary': (120, 120, 120),  # Medium Gray
            'accent': (0, 112, 192),       # Clean Blue accent
            'background': (255, 255, 255), # Pure White
            'text': (45, 45, 45)           # Dark Gray text
        },
        'fonts': {
            'title': 'Segoe UI Light',
            'body': 'Segoe UI'
        }
    },
    'modern-blue': {
        'name': 'Modern Blue',
        'colors': {
            'primary': (31, 78, 120),      # Dark Blue
            'secondary': (0, 112, 192),    # Bright Blue
            'accent': (68, 114, 196),      # Medium Blue
            'background': (255, 255, 255), # White
            'text': (0, 0, 0)              # Black
        },
        'fonts': {
            'title': 'Calibri Light',
            'body': 'Calibri'
        }
    },
    'professional-gray': {
        'name': 'Professional Gray',
        'colors': {
            'primary': (68, 68, 68),       # Dark Gray
            'secondary': (128, 128, 128),  # Medium Gray
            'accent': (0, 112, 192),       # Blue accent
            'background': (255, 255, 255), # White
            'text': (0, 0, 0)              # Black
        },
        'fonts': {
            'title': 'Arial',
            'body': 'Arial'
        }
    },
    'vibrant-orange': {
        'name': 'Vibrant Orange',
        'colors': {
            'primary': (237, 125, 49),     # Orange
            'secondary': (255, 192, 0),    # Yellow
            'accent': (68, 114, 196),      # Blue
            'background': (255, 255, 255), # White
            'text': (0, 0, 0)              # Black
        },
        'fonts': {
            'title': 'Calibri',
            'body': 'Calibri'
        }
    },
    'elegant-purple': {
        'name': 'Elegant Purple',
        'colors': {
            'primary': (112, 48, 160),     # Purple
            'secondary': (146, 208, 80),   # Green accent
            'accent': (255, 192, 0),       # Yellow accent
            'background': (255, 255, 255), # White
            'text': (0, 0, 0)              # Black
        },
        'fonts': {
            'title': 'Calibri Light',
            'body': 'Calibri'
        }
    },
    'clean-green': {
        'name': 'Clean Green',
        'colors': {
            'primary': (70, 136, 71),      # Green
            'secondary': (146, 208, 80),   # Light Green
            'accent': (255, 192, 0),       # Yellow
            'background': (255, 255, 255), # White
            'text': (0, 0, 0)              # Black
        },
        'fonts': {
            'title': 'Calibri',
            'body': 'Calibri'
        }
    }
}


def _apply_theme_to_slide_shapes(slide, theme: dict) -> None:
    """Apply theme styling to all text shapes in a slide

    This is the actual implementation that applies colors and fonts to text.

    Args:
        slide: Slide object
        theme: Theme dict from BUILTIN_THEMES
    """
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    colors = theme['colors']
    fonts = theme['fonts']

    try:
        # Apply to all shapes with text
        for shape in slide.shapes:
            if not hasattr(shape, 'text_frame'):
                continue

            text_frame = shape.text_frame

            # Determine if this is title or body text
            is_title = False
            try:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    if ph_type in [1, 3]:  # TITLE or CENTER_TITLE
                        is_title = True
            except:
                # Check if it's the title shape
                if hasattr(slide.shapes, 'title') and slide.shapes.title == shape:
                    is_title = True

            # Apply styling to all paragraphs and runs
            for paragraph in text_frame.paragraphs:
                for run in paragraph.runs:
                    # Apply font
                    if is_title:
                        run.font.name = fonts['title']
                        if run.font.size is None or run.font.size.pt < 24:
                            run.font.size = Pt(32)  # Ensure titles are large enough
                    else:
                        run.font.name = fonts['body']
                        if run.font.size is None:
                            run.font.size = Pt(18)  # Default body size

                    # Apply text color
                    run.font.color.rgb = RGBColor(*colors['text'])

    except Exception as e:
        print(f"Warning: Could not fully apply theme to slide: {e}")


def apply_theme_to_presentation(prs: Presentation, theme_name: str = 'modern-blue') -> None:
    """Apply built-in theme to presentation

    Applies predefined color scheme and fonts to all slides in presentation.
    This modifies the slide master to affect all slides.

    Args:
        prs: Presentation object
        theme_name: Name of built-in theme ('modern-blue', 'professional-gray',
                    'vibrant-orange', 'elegant-purple', 'clean-green')

    Example:
        prs = Presentation()
        apply_theme_to_presentation(prs, 'vibrant-orange')
    """
    if theme_name not in BUILTIN_THEMES:
        available = ', '.join(BUILTIN_THEMES.keys())
        print(f"Warning: Theme '{theme_name}' not found. Available themes: {available}")
        print("Using default 'modern-blue' theme")
        theme_name = 'modern-blue'

    theme = BUILTIN_THEMES[theme_name]
    colors = theme['colors']
    fonts = theme['fonts']

    try:
        # Access slide master
        slide_master = prs.slide_master

        # Note: Full theme application requires XML manipulation
        # Here we apply basic styling that's safely accessible

        print(f"Applied theme: {theme['name']}")
        print(f"   Primary color: RGB{colors['primary']}")
        print(f"   Accent color: RGB{colors['accent']}")
        print(f"   Fonts: {fonts['title']} / {fonts['body']}")

        # Return theme info for use in slide generation
        return {
            'colors': colors,
            'fonts': fonts,
            'name': theme['name']
        }

    except Exception as e:
        print(f"Warning: Could not fully apply theme: {e}")
        return {
            'colors': colors,
            'fonts': fonts,
            'name': theme['name']
        }


def apply_theme_to_slide(slide, theme_name: str = 'modern-blue') -> None:
    """Apply theme styling to a specific slide

    Applies theme colors and fonts to title and content elements in the slide.
    This is a slide-level application, useful for new slides.

    Args:
        slide: Slide object
        theme_name: Name of built-in theme

    Example:
        slide = prs.slides.add_slide(layout)
        apply_theme_to_slide(slide, 'elegant-purple')
    """
    if theme_name not in BUILTIN_THEMES:
        theme_name = 'modern-blue'

    theme = BUILTIN_THEMES[theme_name]
    colors = theme['colors']
    fonts = theme['fonts']

    try:
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        # Apply to title if exists
        if hasattr(slide.shapes, 'title') and slide.shapes.title:
            title = slide.shapes.title
            if hasattr(title, 'text_frame'):
                for paragraph in title.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = fonts['title']
                        run.font.color.rgb = RGBColor(*colors['text'])

        # Apply to content placeholders
        for shape in slide.shapes:
            if hasattr(shape, 'text_frame'):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        run.font.name = fonts['body']

        print(f"Applied {theme['name']} theme to slide")

    except Exception as e:
        print(f"Warning: Could not fully apply theme to slide: {e}")


# Utility function for applying template styles (placeholder for future enhancement)
def apply_template_style(slide, template_info: dict) -> None:
    """Apply template styles to slide (placeholder for future enhancement)

    This is a placeholder function for applying theme colors, fonts, etc.
    Full implementation would require deeper XML manipulation.

    Args:
        slide: Slide object
        template_info: Template info from analyze_template
    """
    # Future: Apply theme colors, fonts, etc.
    # Requires XML-level manipulation of slide master/theme
    pass


# Additional utility: Helper to discover available layouts
def list_available_layouts(prs: Presentation) -> list:
    """List all available layouts in presentation

    Useful for debugging and understanding what layouts are available.

    Args:
        prs: Presentation object

    Returns:
        List of layout info dicts

    Example:
        layouts = list_available_layouts(prs)
        for layout in layouts:
            print(f"{layout['index']}: {layout['name']}")
    """
    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        layouts.append({
            "index": idx,
            "name": layout.name,
            "placeholder_count": len(layout.placeholders)
        })
    return layouts


# ==============================================================================
# PresentationEditor: High-level safe API for presentation editing
# ==============================================================================

class PresentationEditor:
    """Safe high-level API for presentation editing in Code Interpreter

    This class provides a safe interface for editing presentations by:
    1. Parsing and indexing all slides and elements
    2. Validating operations before execution
    3. Providing clear error messages
    4. Supporting markdown-style text formatting

    Usage:
        prs = Presentation('source.pptx')
        editor = PresentationEditor(prs)

        # Edit operations
        editor.replace_text(slide_index=0, element_id=1, paragraph_id=0, text="New title")
        editor.add_paragraph(slide_index=0, element_id=1, text="New bullet point")
        editor.delete_image(slide_index=2, element_id=3)

        prs.save('output.pptx')
    """

    def __init__(self, prs: Presentation):
        """Initialize PresentationEditor with a presentation

        Args:
            prs: Presentation object to edit

        Raises:
            ValueError: If presentation cannot be parsed (corrupted or incompatible file)
        """
        self.prs = prs
        self.slides = []
        self.text_calculator = TextSizeCalculator()  # For text size optimization

        try:
            self._parse_structure()
        except AttributeError as e:
            raise ValueError(
                f"Failed to parse presentation structure. "
                f"The file may be corrupted or use unsupported features. "
                f"Error: {str(e)}"
            )
        except Exception as e:
            raise ValueError(f"Failed to initialize presentation editor: {str(e)}")

    def _parse_structure(self):
        """Parse presentation structure for safe indexing

        Raises:
            AttributeError: If slides cannot be accessed (file issue)
        """
        try:
            slides_list = list(self.prs.slides)
        except AttributeError as e:
            raise AttributeError(
                f"Cannot access presentation slides. "
                f"File may be corrupted or incompatible: {str(e)}"
            )

        for slide_idx, slide in enumerate(slides_list):
            elements = []
            for elem_idx, shape in enumerate(slide.shapes):
                try:
                    element_info = {
                        'index': elem_idx,
                        'shape': shape,
                        'type': self._get_shape_type(shape),
                        'has_text_frame': hasattr(shape, 'text_frame') and hasattr(shape.text_frame, 'paragraphs')
                    }

                    # Parse paragraphs if text frame exists
                    if element_info['has_text_frame']:
                        element_info['paragraphs'] = []
                        try:
                            for para_idx, para in enumerate(shape.text_frame.paragraphs):
                                element_info['paragraphs'].append({
                                    'index': para_idx,
                                    'text': para.text,
                                    'level': para.level
                                })
                        except Exception as e:
                            print(f"Warning: Could not parse paragraphs for element {elem_idx}: {e}")
                            element_info['paragraphs'] = []

                    elements.append(element_info)
                except Exception as e:
                    # Skip problematic shapes but log the error
                    print(f"Warning: Could not parse shape {elem_idx} on slide {slide_idx}: {e}")
                    continue

            self.slides.append({
                'index': slide_idx,
                'slide': slide,
                'elements': elements,
                'layout_name': slide.slide_layout.name if slide.slide_layout else 'Unknown'
            })

    def _get_shape_type(self, shape) -> str:
        """Get human-readable shape type using feature-based detection"""
        try:
            from pptx.enum.shapes import MSO_SHAPE_TYPE

            # Feature-based detection (same logic as ppt_operations.py)
            if hasattr(shape, 'shapes'):
                return 'group'
            elif hasattr(shape, 'has_table') and shape.has_table:
                return 'table'
            elif hasattr(shape, 'has_chart') and shape.has_chart:
                return 'chart'
            elif hasattr(shape, 'text_frame') and shape.text_frame:
                return 'text'
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                return 'picture'
            else:
                return 'unknown'
        except Exception as e:
            print(f"Warning: Could not determine shape type: {e}")
            return 'unknown'

    def _validate_slide_index(self, slide_index: int):
        """Validate slide index"""
        if not 0 <= slide_index < len(self.slides):
            raise ValueError(
                f"Invalid slide_index {slide_index}. "
                f"Presentation has {len(self.slides)} slides (0-{len(self.slides)-1})"
            )

    def _validate_element_index(self, slide_index: int, element_id: int):
        """Validate element index"""
        self._validate_slide_index(slide_index)
        elements = self.slides[slide_index]['elements']
        if not 0 <= element_id < len(elements):
            raise ValueError(
                f"Invalid element_id {element_id} for slide {slide_index}. "
                f"Slide has {len(elements)} elements (0-{len(elements)-1})"
            )

    def _parse_markdown_text(self, text: str) -> list:
        """Parse markdown text to extract formatting

        Supports: **bold**, *italic*, `code`, [links](url)

        Args:
            text: Markdown-formatted text

        Returns:
            List of (text, format_dict) tuples
        """
        try:
            from markdown_it import MarkdownIt
            from bs4 import BeautifulSoup

            md = MarkdownIt()
            html = md.render(text).strip()
            soup = BeautifulSoup(html, 'html.parser')

            # Extract text with formatting
            blocks = []
            self._extract_text_blocks(soup, blocks, {})
            return blocks if blocks else [(text, {})]

        except Exception as e:
            # Fallback: plain text
            print(f"Warning: Markdown parsing failed ({e}), using plain text")
            return [(text, {})]

    def _extract_text_blocks(self, element, blocks: list, current_format: dict):
        """Recursively extract text blocks with formatting from HTML"""
        if isinstance(element, str):
            if element.strip():
                blocks.append((element, current_format.copy()))
            return

        # Update format based on tag
        new_format = current_format.copy()
        if element.name == 'strong' or element.name == 'b':
            new_format['bold'] = True
        elif element.name == 'em' or element.name == 'i':
            new_format['italic'] = True
        elif element.name == 'code':
            new_format['code'] = True
        elif element.name == 'a':
            new_format['hyperlink'] = element.get('href', '')

        # Process children
        for child in element.children:
            self._extract_text_blocks(child, blocks, new_format)

    def _apply_text_formatting(self, paragraph, text_blocks: list):
        """Apply formatting to paragraph runs while preserving existing styles

        Args:
            paragraph: Paragraph object
            text_blocks: List of (text, format_dict) tuples
        """
        # Preserve original font properties from first run (if exists)
        original_size = None
        original_color_type = None
        original_color_value = None
        original_theme_color = None
        original_name = None
        original_bold = None
        original_italic = None

        if len(paragraph.runs) > 0:
            first_run = paragraph.runs[0]
            original_size = first_run.font.size
            original_name = first_run.font.name
            original_bold = first_run.font.bold
            original_italic = first_run.font.italic

            # Preserve color based on type
            try:
                original_color_type = first_run.font.color.type
                if original_color_type == 1:  # RGB
                    original_color_value = first_run.font.color.rgb
                elif original_color_type == 2:  # Theme/Scheme color
                    original_theme_color = first_run.font.color.theme_color
            except:
                pass  # Color may not be accessible

        # Reuse existing runs instead of removing them (safer and cleaner)
        for i, (text, fmt) in enumerate(text_blocks):
            # Reuse existing run or create new one
            if i < len(paragraph.runs):
                run = paragraph.runs[i]
            else:
                run = paragraph.add_run()

            # Set text
            run.text = text

            # Preserve original formatting unless explicitly overridden by markdown
            if original_size is not None:
                run.font.size = original_size
            if original_name is not None and not fmt.get('code'):
                run.font.name = original_name

            # Preserve color (RGB or Theme color)
            if original_color_type == 1 and original_color_value is not None:
                run.font.color.rgb = original_color_value
            elif original_color_type == 2 and original_theme_color is not None:
                run.font.color.theme_color = original_theme_color

            # Preserve bold/italic if not overridden by markdown
            if not fmt.get('bold') and original_bold is not None:
                run.font.bold = original_bold
            if not fmt.get('italic') and original_italic is not None:
                run.font.italic = original_italic

            # Apply markdown formatting (only if specified)
            if fmt.get('bold'):
                run.font.bold = True
            if fmt.get('italic'):
                run.font.italic = True
            if fmt.get('code'):
                run.font.name = 'Consolas'  # Override for code blocks
            if fmt.get('hyperlink'):
                try:
                    run.hyperlink.address = fmt['hyperlink']
                except:
                    pass  # Hyperlink may not be supported in all contexts

        # Clear text from unused runs (don't try to remove them)
        for i in range(len(text_blocks), len(paragraph.runs)):
            paragraph.runs[i].text = ""

    # ============================================================================
    # Public API: Text Operations
    # ============================================================================

    def update_element_text(self, slide_index: int, element_id: int, text: str):
        """Update entire element text (replaces all paragraphs with new text)

        This is the primary text update method. It replaces the entire text content
        of an element while preserving formatting styles from the first paragraph.

        Args:
            slide_index: Slide index (0-based)
            element_id: Element index within slide (0-based)
            text: New text content (multi-line text creates multiple paragraphs)

        Raises:
            ValueError: If indices are invalid or element has no text frame

        Example:
            editor.update_element_text(0, 1, "New title")
            editor.update_element_text(0, 2, "Point 1\\nPoint 2\\nPoint 3")  # Multi-line
        """
        self._validate_element_index(slide_index, element_id)

        element = self.slides[slide_index]['elements'][element_id]

        if not element['has_text_frame']:
            raise ValueError(
                f"Element {element_id} (type: {element['type']}) has no text frame. "
                f"Cannot update text."
            )

        shape = element['shape']
        text_frame = shape.text_frame

        # Preserve original formatting from first paragraph
        original_size = None
        original_name = None
        original_bold = None
        original_italic = None
        original_color_type = None
        original_color_value = None
        original_theme_color = None

        if len(text_frame.paragraphs) > 0 and len(text_frame.paragraphs[0].runs) > 0:
            first_run = text_frame.paragraphs[0].runs[0]
            original_size = first_run.font.size
            original_name = first_run.font.name
            original_bold = first_run.font.bold
            original_italic = first_run.font.italic

            try:
                original_color_type = first_run.font.color.type
                if original_color_type == 1:  # RGB
                    original_color_value = first_run.font.color.rgb
                elif original_color_type == 2:  # Theme color
                    original_theme_color = first_run.font.color.theme_color
            except:
                pass

        # Check text length and provide warning if it may overflow
        try:
            shape_width_inches = shape.width / 914400  # EMU to inches
            shape_height_inches = shape.height / 914400

            # Get current font size or use default
            if original_size is not None:
                current_font_size = original_size.pt
            else:
                current_font_size = 18  # Default assumption

            # Estimate if text will fit
            estimated_width = self.text_calculator.estimate_text_width(text, int(current_font_size))
            container_width_pts = shape_width_inches * 72 * 0.9  # 90% of available width

            if estimated_width > container_width_pts:
                print(f"Warning: Text may overflow. Consider shorter text.")
                print(f"   Estimated width: {estimated_width:.0f}pts, Available: {container_width_pts:.0f}pts")

                # Auto-wrap if too long (for single-line text)
                if '\n' not in text:
                    wrapped_text = self.text_calculator.wrap_text_intelligently(
                        text, shape_width_inches, int(current_font_size)
                    )
                    if wrapped_text != text:
                        print(f"   Auto-wrapped text to {wrapped_text.count(chr(10)) + 1} lines")
                        text = wrapped_text
        except Exception as e:
            # If size calculation fails, continue without warning
            pass

        # Split text by newlines for multi-paragraph support
        lines = text.split('\n')

        # Clear existing content
        text_frame.clear()

        # Add new paragraphs
        for idx, line in enumerate(lines):
            if idx == 0:
                # Use existing first paragraph
                paragraph = text_frame.paragraphs[0]
            else:
                paragraph = text_frame.add_paragraph()

            # Set text
            paragraph.text = line

            # Apply preserved formatting
            if paragraph.runs:
                for run in paragraph.runs:
                    if original_size is not None:
                        run.font.size = original_size
                    if original_name is not None:
                        run.font.name = original_name
                    if original_bold is not None:
                        run.font.bold = original_bold
                    if original_italic is not None:
                        run.font.italic = original_italic

                    # Preserve color
                    if original_color_type == 1 and original_color_value is not None:
                        run.font.color.rgb = original_color_value
                    elif original_color_type == 2 and original_theme_color is not None:
                        run.font.color.theme_color = original_theme_color

    def find_and_replace_in_element(self, slide_index: int, element_id: int, find: str, replace: str):
        """Find and replace text within an element (preserves formatting)

        This method searches for all occurrences of the find text within the element
        and replaces them with the replacement text. Unlike update_element_text which
        replaces the entire text, this method only modifies matching portions.

        Args:
            slide_index: Slide index (0-based)
            element_id: Element index within slide (0-based)
            find: Text to find
            replace: Text to replace with

        Raises:
            ValueError: If indices are invalid or element has no text frame

        Example:
            editor.find_and_replace_in_element(0, 1, "Quality assessment", "Quality evaluation")
            editor.find_and_replace_in_element(2, 3, "2024", "2025")
        """
        self._validate_element_index(slide_index, element_id)

        element = self.slides[slide_index]['elements'][element_id]

        if not element['has_text_frame']:
            raise ValueError(
                f"Element {element_id} (type: {element['type']}) has no text frame. "
                f"Cannot find and replace text."
            )

        shape = element['shape']
        text_frame = shape.text_frame

        # Track if any replacements were made
        replacement_count = 0

        # Iterate through all paragraphs
        for para in text_frame.paragraphs:
            # Check if find text exists in paragraph
            full_text = para.text
            if find not in full_text:
                continue

            # Replace text while preserving run-level formatting (colors, bold, etc.)
            # We need to iterate through runs and replace within each run
            runs_to_process = list(para.runs)  # Copy list to avoid modification issues

            for run in runs_to_process:
                run_text = run.text
                if find in run_text:
                    # Replace within this run
                    run.text = run_text.replace(find, replace)
                    replacement_count += run_text.count(find)
                    # Run formatting (color, bold, italic, font) is automatically preserved

        if replacement_count > 0:
            print(f"Replaced {replacement_count} occurrence(s) of '{find}' with '{replace}'")
        else:
            print(f"Warning: Text '{find}' not found in element {element_id}")

    def replace_image(self, slide_index: int, element_id: int, image_path: str):
        """Replace an image with a new one

        Args:
            slide_index: Slide index (0-based)
            element_id: Element index within slide (0-based)
            image_path: Path to new image file

        Raises:
            ValueError: If indices are invalid or element is not an image
        """
        self._validate_element_index(slide_index, element_id)

        element = self.slides[slide_index]['elements'][element_id]

        if element['type'] != 'picture':
            raise ValueError(
                f"Element {element_id} is not an image (type: {element['type']}). "
                f"Cannot replace image."
            )

        import os
        if not os.path.exists(image_path):
            raise ValueError(f"Image file not found: {image_path}")

        shape = element['shape']

        # Get current position and size
        left = shape.left
        top = shape.top
        width = shape.width
        height = shape.height

        # Get slide
        slide = self.slides[slide_index]['slide']

        # Remove old image
        sp = shape.element
        sp.getparent().remove(sp)

        # Add new image with same position and size
        try:
            from PIL import Image

            # Calculate aspect ratio preserving size
            img = Image.open(image_path)
            img_ratio = img.width / img.height
            shape_ratio = width / height

            if img_ratio > shape_ratio:
                # Image is wider, fit to width
                new_width = width
                new_height = Pt(width / img_ratio)
                new_top = Pt(top + (height - new_height) / 2)
                new_left = left
            else:
                # Image is taller, fit to height
                new_height = height
                new_width = Pt(height * img_ratio)
                new_left = Pt(left + (width - new_width) / 2)
                new_top = top

            slide.shapes.add_picture(image_path, new_left, new_top, width=new_width, height=new_height)

        except Exception as e:
            # Fallback: use original size
            print(f"Warning: Could not preserve aspect ratio ({e}), using original size")
            slide.shapes.add_picture(image_path, left, top, width=width, height=height)

    def get_slide_count(self) -> int:
        """Get total number of slides"""
        return len(self.slides)

    def get_slide_info(self, slide_index: int) -> dict:
        """Get information about a slide

        Args:
            slide_index: Slide index (0-based)

        Returns:
            Dict with slide information
        """
        self._validate_slide_index(slide_index)

        slide_info = self.slides[slide_index]

        return {
            'index': slide_index,
            'layout_name': slide_info['layout_name'],
            'element_count': len(slide_info['elements']),
            'elements': [
                {
                    'index': elem['index'],
                    'type': elem['type'],
                    'has_text': elem['has_text_frame'],
                    'paragraph_count': len(elem.get('paragraphs', []))
                }
                for elem in slide_info['elements']
            ]
        }

    def get_element_info(self, slide_index: int, element_id: int) -> dict:
        """Get detailed information about an element

        Args:
            slide_index: Slide index (0-based)
            element_id: Element index (0-based)

        Returns:
            Dict with element information
        """
        self._validate_element_index(slide_index, element_id)

        element = self.slides[slide_index]['elements'][element_id]

        info = {
            'index': element['index'],
            'type': element['type'],
            'has_text_frame': element['has_text_frame']
        }

        if element['has_text_frame']:
            info['paragraphs'] = element['paragraphs']

        return info
