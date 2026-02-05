"""
PowerPoint Operations - Convert high-level operations to safe python-pptx code

This module converts structured operations (dicts) into safe, executable
python-pptx code for Code Interpreter execution.

The approach:
1. Agent decides WHAT to do (high-level operations)
2. Tool converts to HOW (safe python-pptx code)
3. Code Interpreter executes (sandboxed environment)

This ensures:
- Agent never writes low-level code directly
- All code is validated before execution
- Index management is automatic and safe
- Clear error messages for debugging

Element Detection Strategy:
- Feature-based detection (hasattr) instead of shape_type enum
- This catches ALL text-editable shapes including animations, custom shapes
- Types: text, picture, table, chart, group, unknown
- Role information (TITLE, BODY) preserved for context
"""

from typing import Dict, List, Any


def generate_batch_update_slides_code(
    source_filename: str,
    output_filename: str,
    slide_updates: List[Dict[str, Any]]
) -> str:
    """Generate safe python-pptx code for updating multiple slides

    Args:
        source_filename: Source presentation filename
        output_filename: Output presentation filename
        slide_updates: List of dicts with 'slide_index' and 'operations'

    Returns:
        Complete python code string ready for execution

    Example slide_updates:
        [
            {"slide_index": 0, "operations": [{"action": "set_text", ...}]},
            {"slide_index": 2, "operations": [{"action": "replace_text", ...}]}
        ]
    """
    # Build code for each slide
    all_slide_code = []

    for slide_idx, update in enumerate(slide_updates):
        slide_index = update['slide_index']
        operations = update['operations']

        # Build operation code lines for this slide
        operation_lines = []

        for op_idx, op in enumerate(operations):
            action = op.get('action')

            if not action:
                raise ValueError(f"Slide {slide_index}, operation {op_idx} missing 'action' field")

            try:
                if action == 'set_text':
                    code_line = _generate_set_text(slide_index, op)
                elif action == 'replace_text':
                    code_line = _generate_replace_text(slide_index, op)
                elif action == 'replace_image':
                    code_line = _generate_replace_image(slide_index, op)
                else:
                    raise ValueError(f"Unknown action '{action}' in operation {op_idx}. Supported: 'set_text', 'replace_text', 'replace_image'")

                operation_lines.append(f"# Operation {op_idx + 1}: {action}")
                operation_lines.append(code_line)

            except KeyError as e:
                raise ValueError(
                    f"Slide {slide_index}, operation {op_idx} (action: {action}) missing required field: {e}"
                )
            except Exception as e:
                raise ValueError(
                    f"Error generating code for slide {slide_index}, operation {op_idx} (action: {action}): {e}"
                )

        # Combine operations for this slide
        operations_code = '\n'.join(operation_lines)
        slide_code = f"""
# === Slide {slide_index + 1} ({len(operations)} operation(s)) ===
{operations_code}
"""
        all_slide_code.append(slide_code)

    # Combine all slide code
    all_slides_combined = '\n'.join(all_slide_code)

    # Generate final code
    total_operations = sum(len(update['operations']) for update in slide_updates)
    code = f"""
from pptx import Presentation
from presentation_editor import PresentationEditor

# Load presentation
prs = Presentation('{source_filename}')
editor = PresentationEditor(prs)

# Execute operations on {len(slide_updates)} slide(s)
{all_slides_combined}

# Save result
prs.save('{output_filename}')

# Print summary
print(f"Successfully updated {len(slide_updates)} slide(s)")
print(f"Applied {total_operations} total operation(s)")
""".strip()

    return code


def _generate_set_text(slide_index: int, op: dict) -> str:
    """Generate code for set_text operation - replaces entire element text

    Args:
        slide_index: Slide index (0-based)
        op: Operation dict with 'element_id' and 'text'

    Returns:
        Python code string for execution
    """
    element_id = op['element_id']
    text = op['text']

    # Use repr() to safely escape all special characters
    text_repr = repr(text)

    return f"editor.update_element_text({slide_index}, {element_id}, {text_repr})"


def _generate_replace_text(slide_index: int, op: dict) -> str:
    """Generate code for replace_text operation - find and replace text within element

    Args:
        slide_index: Slide index (0-based)
        op: Operation dict with 'element_id', 'find', and 'replace'

    Returns:
        Python code string for execution
    """
    element_id = op['element_id']
    find_text = op['find']
    replace_text = op['replace']

    # Use repr() to safely escape all special characters
    find_repr = repr(find_text)
    replace_repr = repr(replace_text)

    return f"editor.find_and_replace_in_element({slide_index}, {element_id}, {find_repr}, {replace_repr})"


def _generate_replace_image(slide_index: int, op: dict) -> str:
    """Generate code for replace_image operation - replaces entire element image

    Args:
        slide_index: Slide index (0-based)
        op: Operation dict with 'element_id' and 'image_path'

    Returns:
        Python code string for execution
    """
    element_id = op['element_id']
    image_path = op['image_path']

    return f"editor.replace_image({slide_index}, {element_id}, '{image_path}')"


def generate_add_slide_code(
    source_filename: str,
    output_filename: str,
    layout_name: str,
    position: int,
    custom_code: str = None
) -> str:
    """Generate code for adding a new slide with optional custom python-pptx code

    Args:
        source_filename: Source presentation filename
        output_filename: Output presentation filename
        layout_name: Name of layout to use
        position: Position to insert (-1 for append)
        custom_code: Optional Python code to execute on the new slide
                     Available variables: slide, prs, Inches, Pt, RGBColor

    Returns:
        Python code string
    """
    # Prepare custom code section
    if custom_code:
        # Indent the custom code properly (4 spaces for base level)
        indented_custom = '\n'.join('    ' + line if line.strip() else ''
                                     for line in custom_code.split('\n'))
        custom_section = f"""
    # Execute custom code
    slide = new_slide
{indented_custom}
"""
    else:
        custom_section = ""

    code = f"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.chart.data import CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

try:
    # Load presentation
    prs = Presentation('{source_filename}')

    # Verify slides are accessible (catch corrupted files early)
    try:
        slide_count = len(prs.slides)
    except AttributeError as e:
        raise ValueError(f"Cannot access slides in presentation. File may be corrupted: {{str(e)}}")

    # Find layout by name
    layout = None
    try:
        for l in prs.slide_layouts:
            if l.name == '{layout_name}':
                layout = l
                break
    except AttributeError as e:
        raise ValueError(f"Cannot access slide layouts. File may be corrupted: {{str(e)}}")

    if layout is None:
        available_layouts = []
        try:
            available_layouts = [l.name for l in prs.slide_layouts]
        except:
            available_layouts = ["<unable to list layouts>"]
        raise ValueError(f"Layout '{layout_name}' not found. Available layouts: {{available_layouts}}")

    # Add slide
    try:
        new_slide = prs.slides.add_slide(layout)
    except AttributeError as e:
        raise ValueError(f"Cannot add slide to presentation. File may be corrupted: {{str(e)}}")
{custom_section}
    # Move to position if specified
    if {position} >= 0 and {position} < len(prs.slides):
        # Work with slide ID list (XML elements), not slide objects
        xml_slides = prs.slides._sldIdLst
        slides_list = list(xml_slides)

        # Move newly added slide (last position) to target position
        new_slide_id = slides_list.pop()  # Remove from end
        slides_list.insert({position}, new_slide_id)  # Insert at target position

        # Rebuild slide order
        xml_slides.clear()
        for slide_id in slides_list:
            xml_slides.append(slide_id)

    # Save
    prs.save('{output_filename}')

    print(f"Added slide with layout '{{layout.name}}' at position {{{position} if {position} >= 0 else 'end'}}")

except ValueError as e:
    # Re-raise with original message
    raise
except Exception as e:
    raise ValueError(f"Failed to add slide: {{str(e)}}")
""".strip()

    return code


def generate_delete_slides_code(
    source_filename: str,
    output_filename: str,
    slide_indices: List[int]
) -> str:
    """Generate code for deleting multiple slides

    Args:
        source_filename: Source presentation filename
        output_filename: Output presentation filename
        slide_indices: List of slide indices to delete (0-based)

    Returns:
        Python code string
    """
    # Sort indices in reverse order to delete from end to beginning
    # This prevents index shifting issues
    sorted_indices = sorted(set(slide_indices), reverse=True)

    code = f"""
from pptx import Presentation

try:
    # Load presentation
    prs = Presentation('{source_filename}')

    # Verify slides are accessible
    try:
        total_slides = len(prs.slides)
    except AttributeError as e:
        raise ValueError(f"Cannot access slides. File may be corrupted: {{str(e)}}")

    # Validate indices
    indices_to_delete = {sorted_indices}

    for idx in indices_to_delete:
        if not 0 <= idx < total_slides:
            raise ValueError(f"Invalid slide index {{idx}}. Presentation has {{total_slides}} slides (0-{{total_slides-1}})")

    # Delete slides (from end to beginning to avoid index shifts)
    for idx in indices_to_delete:
        rId = prs.slides._sldIdLst[idx].rId
        prs.part.drop_rel(rId)
        del prs.slides._sldIdLst[idx]

    # Save
    prs.save('{output_filename}')

    print(f"Deleted {{len(indices_to_delete)}} slide(s)")
    print(f"Remaining slides: {{len(prs.slides)}}")

except ValueError as e:
    raise
except Exception as e:
    raise ValueError(f"Failed to delete slides: {{str(e)}}")
""".strip()

    return code


def generate_move_slide_code(
    source_filename: str,
    output_filename: str,
    from_index: int,
    to_index: int
) -> str:
    """Generate code for moving a slide

    Args:
        source_filename: Source presentation filename
        output_filename: Output presentation filename
        from_index: Source position (0-based)
        to_index: Target position (0-based)

    Returns:
        Python code string
    """
    code = f"""
from pptx import Presentation

try:
    # Load presentation
    prs = Presentation('{source_filename}')

    # Verify slides are accessible
    try:
        total_slides = len(prs.slides)
    except AttributeError as e:
        raise ValueError(f"Cannot access slides. File may be corrupted: {{str(e)}}")

    # Validate indices
    if not 0 <= {from_index} < total_slides:
        raise ValueError(f"Invalid from_index {{{from_index}}}. Presentation has {{total_slides}} slides")
    if not 0 <= {to_index} < total_slides:
        raise ValueError(f"Invalid to_index {{{to_index}}}. Presentation has {{total_slides}} slides")

    # Move slide
    xml_slides = prs.slides._sldIdLst
    slides_list = list(xml_slides)

    # Remove from source position and insert at target position
    slide_to_move = slides_list.pop({from_index})
    slides_list.insert({to_index}, slide_to_move)

    # Rebuild slide order
    xml_slides.clear()
    for slide in slides_list:
        xml_slides.append(slide)

    # Save
    prs.save('{output_filename}')

    print(f"Moved slide from position {{{from_index} + 1}} to {{{to_index} + 1}}")

except ValueError as e:
    raise
except Exception as e:
    raise ValueError(f"Failed to move slide: {{str(e)}}")
""".strip()

    return code


def generate_duplicate_slide_code(
    source_filename: str,
    output_filename: str,
    slide_index: int,
    position: int
) -> str:
    """Generate code for duplicating a slide

    Args:
        source_filename: Source presentation filename
        output_filename: Output presentation filename
        slide_index: Index of slide to duplicate (0-based)
        position: Position for duplicate (-1 for after original)

    Returns:
        Python code string
    """
    code = f"""
from pptx import Presentation
from copy import deepcopy

try:
    # Load presentation
    prs = Presentation('{source_filename}')

    # Verify slides are accessible
    try:
        total_slides = len(prs.slides)
    except AttributeError as e:
        raise ValueError(f"Cannot access slides. File may be corrupted: {{str(e)}}")

    # Validate index
    if not 0 <= {slide_index} < total_slides:
        raise ValueError(f"Invalid slide_index {{{slide_index}}}. Presentation has {{total_slides}} slides")

    # Get source slide
    source_slide = prs.slides[{slide_index}]

    # Duplicate slide by adding new slide with same layout, then copying content
    layout = source_slide.slide_layout
    new_slide = prs.slides.add_slide(layout)

    # Copy all shapes from source to new slide
    for shape in source_slide.shapes:
        try:
            # Get shape element and create deep copy
            from pptx.oxml import parse_xml
            shape_xml = shape.element.xml
            new_shape_element = parse_xml(shape_xml)
            new_slide.shapes._spTree.insert_element_before(new_shape_element, 'p:extLst')
        except Exception as e:
            print(f"Warning: Could not copy shape: {{str(e)}}")
            continue

    # Copy speaker notes if they exist
    try:
        if source_slide.has_notes_slide:
            source_notes = source_slide.notes_slide.notes_text_frame.text
            if source_notes:
                new_notes = new_slide.notes_slide
                new_notes.notes_text_frame.text = source_notes
    except:
        pass

    # Move new slide to target position if specified
    insert_pos = {position} if {position} >= 0 else {slide_index} + 1
    if insert_pos < len(prs.slides) - 1:
        # Work with slide ID list
        xml_slides = prs.slides._sldIdLst
        slides_list = list(xml_slides)

        # Move newly added slide (last position) to target position
        new_slide_id = slides_list.pop()  # Remove from end
        slides_list.insert(insert_pos, new_slide_id)  # Insert at target

        # Rebuild slide order
        xml_slides.clear()
        for slide_id in slides_list:
            xml_slides.append(slide_id)

    # Save
    prs.save('{output_filename}')

    print(f"Duplicated slide {{{slide_index} + 1}} to position {{insert_pos + 1}}")

except ValueError as e:
    raise
except Exception as e:
    raise ValueError(f"Failed to duplicate slide: {{str(e)}}")
""".strip()

    return code


def generate_update_notes_code(
    source_filename: str,
    output_filename: str,
    slide_index: int,
    notes_text: str
) -> str:
    """Generate code for updating slide notes

    Args:
        source_filename: Source presentation filename
        output_filename: Output presentation filename
        slide_index: Slide index (0-based)
        notes_text: New notes content

    Returns:
        Python code string
    """
    # Use repr() for safe string escaping
    notes_repr = repr(notes_text)

    code = f"""
from pptx import Presentation

try:
    # Load presentation
    prs = Presentation('{source_filename}')

    # Verify slides are accessible
    try:
        total_slides = len(prs.slides)
    except AttributeError as e:
        raise ValueError(f"Cannot access slides. File may be corrupted: {{str(e)}}")

    # Validate index
    if not 0 <= {slide_index} < total_slides:
        raise ValueError(f"Invalid slide_index {{{slide_index}}}. Presentation has {{total_slides}} slides (0-{{total_slides-1}})")

    # Get slide
    slide = prs.slides[{slide_index}]

    # Update notes
    try:
        notes_slide = slide.notes_slide
        text_frame = notes_slide.notes_text_frame
        text_frame.clear()  # Clear existing notes
        text_frame.text = {notes_repr}
    except Exception as e:
        raise ValueError(f"Failed to update notes: {{str(e)}}")

    # Save
    prs.save('{output_filename}')

    print(f"Updated notes for slide {{{slide_index} + 1}}")

except ValueError as e:
    raise
except Exception as e:
    raise ValueError(f"Failed to update slide notes: {{str(e)}}")
""".strip()

    return code


def generate_analyze_presentation_code(
    presentation_filename: str,
    slide_index: int = None
) -> str:
    """Generate code for analyzing presentation structure

    Args:
        presentation_filename: Presentation filename
        slide_index: Optional slide index (0-based). If provided, only analyzes that slide.

    Returns:
        Python code string that outputs JSON with full content
    """
    # Add slide filter if specific slide requested
    slide_filter = ""
    if slide_index is not None:
        slide_filter = f"""
        # Filter to specific slide only
        if slide_idx != {slide_index}:
            continue
"""

    code = f"""
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
import json

try:
    # Load presentation
    prs = Presentation('{presentation_filename}')

    # Verify slides and layouts are accessible
    try:
        _ = len(prs.slides)
        _ = len(prs.slide_layouts)
    except AttributeError as e:
        raise ValueError(f"Cannot access presentation structure. File may be corrupted: {{str(e)}}")

    # Extract layouts
    layouts = {{}}
    for idx, layout in enumerate(prs.slide_layouts):
        layouts[layout.name] = {{
            'index': idx,
            'name': layout.name,
            'placeholder_count': len(layout.placeholders)
        }}

    # Extract slides with full content (excluding hidden slides)
    slides = []
    for slide_idx, slide in enumerate(prs.slides):
{slide_filter}
        # Skip hidden slides
        try:
            # In PowerPoint XML, show="0" means hidden
            if slide.element.get('show') == '0':
                continue
        except:
            pass  # If attribute doesn't exist, slide is visible

        slide_info = {{
            'index': slide_idx,
            'layout': slide.slide_layout.name if slide.slide_layout else 'Unknown',
            'title': slide.shapes.title.text if slide.shapes.title else None
        }}

        # Extract speaker notes
        try:
            if slide.has_notes_slide:
                notes_text = slide.notes_slide.notes_text_frame.text.strip()
                if notes_text:
                    slide_info['notes'] = notes_text
        except:
            pass

        # Extract elements
        elements = []
        for elem_idx, shape in enumerate(slide.shapes):
            element_info = {{
                'element_id': elem_idx,
                'type': 'unknown',
                'role': None
            }}

            # Extract position and size (EMU to inches: 914400 EMU = 1 inch)
            try:
                element_info['position'] = {{
                    'left': round(shape.left / 914400, 2),
                    'top': round(shape.top / 914400, 2),
                    'width': round(shape.width / 914400, 2),
                    'height': round(shape.height / 914400, 2)
                }}
            except:
                pass

            # Check if placeholder and get role
            try:
                if shape.is_placeholder:
                    ph_type = shape.placeholder_format.type
                    # Map placeholder types to roles
                    if ph_type == 1:
                        element_info['role'] = 'title'
                    elif ph_type == 2:
                        element_info['role'] = 'body'
                    elif ph_type == 4:
                        element_info['role'] = 'subtitle'
                    elif ph_type == 8:
                        element_info['role'] = 'footer'
            except:
                pass

            # Determine type using feature-based detection (not shape_type enum)
            # This ensures we catch all text-editable shapes including animations

            if hasattr(shape, 'shapes'):
                # Group shape
                element_info['type'] = 'group'

            elif hasattr(shape, 'has_table') and shape.has_table:
                # Table shape (GraphicFrame with table)
                element_info['type'] = 'table'
                element_info['rows'] = len(shape.table.rows)
                element_info['cols'] = len(shape.table.columns)

                # Extract table cell text content
                table_data = []
                table = shape.table
                for row_idx, row in enumerate(table.rows):
                    row_cells = []
                    for col_idx, cell in enumerate(row.cells):
                        try:
                            cell_text = cell.text_frame.text.strip() if cell.text_frame else ""
                            row_cells.append(cell_text)
                        except:
                            row_cells.append("")
                    table_data.append(row_cells)
                element_info['table_data'] = table_data

            elif hasattr(shape, 'has_chart') and shape.has_chart:
                # Chart shape (GraphicFrame with chart)
                element_info['type'] = 'chart'

                # Extract chart data
                try:
                    chart = shape.chart
                    chart_info = {{
                        'chart_type': str(chart.chart_type)
                    }}

                    # Extract categories and series data
                    if chart.plots:
                        plot = chart.plots[0]

                        # Categories (X-axis labels)
                        try:
                            if plot.categories:
                                chart_info['categories'] = [str(cat) for cat in plot.categories]
                        except:
                            pass

                        # Series data
                        try:
                            series_data = []
                            for series in chart.series:
                                series_info = {{
                                    'name': series.name if hasattr(series, 'name') else 'Series'
                                }}
                                try:
                                    series_info['values'] = [float(v) if v is not None else 0 for v in series.values]
                                except:
                                    pass
                                series_data.append(series_info)

                            if series_data:
                                chart_info['series'] = series_data
                        except:
                            pass

                    element_info['chart_data'] = chart_info
                except Exception as e:
                    # If chart extraction fails, mark as chart but without data
                    element_info['chart_error'] = str(e)[:100]

            elif hasattr(shape, 'text_frame') and shape.text_frame:
                # Text-editable shape (textbox, placeholder, autoshape, animation text, etc.)
                element_info['type'] = 'text'

            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Picture shape (only check if not text-editable)
                element_info['type'] = 'picture'

            else:
                # Unknown type
                element_info['type'] = 'unknown'

            # Extract text content for text type (has text_frame)
            if element_info['type'] == 'text' and hasattr(shape, 'text_frame'):
                paragraphs = []
                try:
                    for para_idx, para in enumerate(shape.text_frame.paragraphs):
                        # Include ALL paragraphs (even empty ones) - they're still editable
                        para_info = {{
                            'paragraph_id': para_idx,
                            'text': para.text,
                            'level': para.level
                        }}

                        # Extract format info from first run (representative)
                        if len(para.runs) > 0:
                            first_run = para.runs[0]
                            format_info = {{}}

                            try:
                                if first_run.font.name:
                                    format_info['font_name'] = first_run.font.name
                            except:
                                pass

                            try:
                                if first_run.font.size:
                                    format_info['font_size'] = int(first_run.font.size.pt)
                            except:
                                pass

                            try:
                                if first_run.font.bold is not None:
                                    format_info['bold'] = first_run.font.bold
                            except:
                                pass

                            try:
                                if first_run.font.italic is not None:
                                    format_info['italic'] = first_run.font.italic
                            except:
                                pass

                            try:
                                color = first_run.font.color
                                if color.type == 1:  # RGB
                                    # color.rgb returns an int in 0xRRGGBB format
                                    rgb_int = int(color.rgb)
                                    r = (rgb_int >> 16) & 0xFF
                                    g = (rgb_int >> 8) & 0xFF
                                    b = rgb_int & 0xFF
                                    format_info['color'] = f"RGB({{r}},{{g}},{{b}})"
                                elif color.type == 2:  # Theme color
                                    format_info['color'] = f"Theme({{color.theme_color}})"
                            except:
                                pass

                            if format_info:
                                para_info['format'] = format_info

                        paragraphs.append(para_info)
                except:
                    # If text_frame access fails, leave paragraphs empty
                    pass
                element_info['paragraphs'] = paragraphs

            elements.append(element_info)

        slide_info['elements'] = elements
        slides.append(slide_info)

    # Build result
    result = {{
        'total_slides': len(prs.slides),
        'layouts': layouts,
        'slides': slides
    }}

    # Output as JSON
    print(json.dumps(result, indent=2, ensure_ascii=False))

except ValueError as e:
    raise
except Exception as e:
    raise ValueError(f"Failed to analyze presentation: {{str(e)}}")
""".strip()

    return code
