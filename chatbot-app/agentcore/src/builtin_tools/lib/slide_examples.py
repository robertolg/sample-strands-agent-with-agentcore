"""
Slide Code Examples - Reference patterns for python-pptx slide creation.

These are reference examples to adapt for your content needs.
Also useful for debugging code errors.

Categories:
- text_layout: Text arrangement, hierarchy, bullets
- number_highlight: Numbers/KPI emphasis
- grid_layout: Multiple items, comparisons
- image_text: Image + text combinations
- visual_emphasis: Highlights, color boxes, accents
"""

SLIDE_EXAMPLES = {
    "text_layout": {
        "description": "Text arrangement, hierarchy, bullet points",
        "when_to_use": "Multiple text items, lists, structured content",
        "examples": [
            {
                "name": "title_with_bullets",
                "code": '''
# Title with bullet points
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Background
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.8))
tf = title.text_frame
tf.paragraphs[0].text = "Key Findings"
tf.paragraphs[0].font.size = Pt(32)
tf.paragraphs[0].font.bold = True
tf.paragraphs[0].font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Bullet points
bullets = [
    "Revenue increased by 25% year-over-year",
    "Customer satisfaction reached 4.8/5.0",
    "Market share expanded to 35%",
    "Operating costs reduced by 12%"
]

content = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(9), Inches(5))
tf = content.text_frame
tf.word_wrap = True

for i, bullet in enumerate(bullets):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.text = f"• {bullet}"
    p.font.size = Pt(20)
    p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
    p.space_after = Pt(12)
'''
            },
            {
                "name": "hierarchy_text",
                "code": '''
# Text with visual hierarchy (title > subtitle > body)
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xF5, 0xF5, 0xF5)

# Main title - largest
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.7))
p = title.text_frame.paragraphs[0]
p.text = "Digital Transformation"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

# Subtitle - medium
subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
p = subtitle.text_frame.paragraphs[0]
p.text = "Modernizing our technology stack for the future"
p.font.size = Pt(20)
p.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)

# Body text - smallest
body = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(4))
tf = body.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Our three-year roadmap focuses on cloud migration, AI integration, and process automation. This initiative will streamline operations and enhance customer experience."
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
p.line_spacing = 1.5
'''
            }
        ]
    },

    "number_highlight": {
        "description": "Numbers and KPI emphasis with large fonts",
        "when_to_use": "Key metrics, statistics, performance data",
        "examples": [
            {
                "name": "single_big_number",
                "code": '''
# Single big number with label
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Big number - centered
number = slide.shapes.add_textbox(Inches(0), Inches(2), Inches(10), Inches(2))
tf = number.text_frame
p = tf.paragraphs[0]
p.text = "147%"
p.font.size = Pt(120)
p.font.bold = True
p.font.color.rgb = RGBColor(0x00, 0xD4, 0xAA)
p.alignment = PP_ALIGN.CENTER

# Label below
label = slide.shapes.add_textbox(Inches(0), Inches(4.2), Inches(10), Inches(0.8))
tf = label.text_frame
p = tf.paragraphs[0]
p.text = "Revenue Growth"
p.font.size = Pt(28)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER
'''
            },
            {
                "name": "multiple_metrics",
                "code": '''
# Multiple KPIs in a row
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
p = title.text_frame.paragraphs[0]
p.text = "Q4 Performance"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

# KPI data
metrics = [
    ("$2.4M", "Revenue"),
    ("1,250", "New Customers"),
    ("98.5%", "Uptime"),
]

# Create KPI boxes
start_x = 0.5
box_width = 3
spacing = 0.25

for i, (value, label) in enumerate(metrics):
    x = start_x + i * (box_width + spacing)

    # Value
    val_box = slide.shapes.add_textbox(Inches(x), Inches(2.5), Inches(box_width), Inches(1.2))
    p = val_box.text_frame.paragraphs[0]
    p.text = value
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x7A, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    # Label
    lbl_box = slide.shapes.add_textbox(Inches(x), Inches(3.7), Inches(box_width), Inches(0.5))
    p = lbl_box.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(18)
    p.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
    p.alignment = PP_ALIGN.CENTER
'''
            }
        ]
    },

    "grid_layout": {
        "description": "Multiple items in grid arrangement",
        "when_to_use": "Comparisons, feature lists, team members, product lineup",
        "examples": [
            {
                "name": "two_column",
                "code": '''
# Two column comparison
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
p = title.text_frame.paragraphs[0]
p.text = "Before vs After"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

columns = [
    {"title": "Before", "color": RGBColor(0xE7, 0x4C, 0x3C), "items": ["Manual processes", "3-day turnaround", "High error rate"]},
    {"title": "After", "color": RGBColor(0x27, 0xAE, 0x60), "items": ["Automated workflow", "Same-day delivery", "99.9% accuracy"]}
]

for i, col in enumerate(columns):
    x = 0.5 + i * 4.75

    # Column header
    header = slide.shapes.add_textbox(Inches(x), Inches(1.5), Inches(4.5), Inches(0.6))
    p = header.text_frame.paragraphs[0]
    p.text = col["title"]
    p.font.size = Pt(24)
    p.font.bold = True
    p.font.color.rgb = col["color"]

    # Column items
    content = slide.shapes.add_textbox(Inches(x), Inches(2.2), Inches(4.5), Inches(4))
    tf = content.text_frame
    tf.word_wrap = True

    for j, item in enumerate(col["items"]):
        if j == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
        p.space_after = Pt(8)
'''
            },
            {
                "name": "three_column_cards",
                "code": '''
# Three column card layout
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xF0, 0xF0, 0xF0)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
p = title.text_frame.paragraphs[0]
p.text = "Our Services"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

cards = [
    {"title": "Consulting", "desc": "Expert guidance for digital transformation"},
    {"title": "Development", "desc": "Custom software solutions"},
    {"title": "Support", "desc": "24/7 technical assistance"}
]

card_width = 2.9
start_x = 0.5
spacing = 0.3

for i, card in enumerate(cards):
    x = start_x + i * (card_width + spacing)

    # Card background (white rectangle)
    bg = slide.shapes.add_shape(1, Inches(x), Inches(1.2), Inches(card_width), Inches(4.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
    bg.line.fill.background()

    # Card title
    t = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.5), Inches(card_width - 0.4), Inches(0.6))
    p = t.text_frame.paragraphs[0]
    p.text = card["title"]
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

    # Card description
    d = slide.shapes.add_textbox(Inches(x + 0.2), Inches(2.2), Inches(card_width - 0.4), Inches(3))
    tf = d.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = card["desc"]
    p.font.size = Pt(14)
    p.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
'''
            }
        ]
    },

    "image_text": {
        "description": "Image and text combinations",
        "when_to_use": "Product showcase, team intro, visual storytelling",
        "examples": [
            {
                "name": "left_image_right_text",
                "code": '''
# Left image, right text layout
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import os

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Image placeholder (left side) - use actual image if available
# Check for images in workspace
image_files = [f for f in os.listdir('.') if f.endswith(('.png', '.jpg', '.jpeg'))]
if image_files:
    slide.shapes.add_picture(image_files[0], Inches(0.5), Inches(0.5), width=Inches(4.5))
else:
    # Placeholder box if no image
    placeholder = slide.shapes.add_shape(1, Inches(0.5), Inches(0.5), Inches(4.5), Inches(6))
    placeholder.fill.solid()
    placeholder.fill.fore_color.rgb = RGBColor(0xE0, 0xE0, 0xE0)
    placeholder.line.fill.background()

# Text content (right side)
title = slide.shapes.add_textbox(Inches(5.3), Inches(1), Inches(4.2), Inches(0.8))
p = title.text_frame.paragraphs[0]
p.text = "Product Name"
p.font.size = Pt(28)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

desc = slide.shapes.add_textbox(Inches(5.3), Inches(2), Inches(4.2), Inches(4))
tf = desc.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Description of the product or feature. Explain the key benefits and value proposition here."
p.font.size = Pt(16)
p.font.color.rgb = RGBColor(0x33, 0x33, 0x33)
p.line_spacing = 1.4
'''
            },
            {
                "name": "full_background_overlay",
                "code": '''
# Full background with text overlay
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
import os

# Dark overlay background (simulates image with overlay)
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# If image available, add it first (will be behind text)
image_files = [f for f in os.listdir('.') if f.endswith(('.png', '.jpg', '.jpeg'))]
if image_files:
    slide.shapes.add_picture(image_files[0], Inches(0), Inches(0), width=Inches(10), height=Inches(7.5))
    # Semi-transparent overlay
    overlay = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(10), Inches(7.5))
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = RGBColor(0x00, 0x00, 0x00)
    # Note: transparency requires XML manipulation, using dark color instead

# Centered text
title = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(9), Inches(1.5))
p = title.text_frame.paragraphs[0]
p.text = "Innovation Starts Here"
p.font.size = Pt(48)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.alignment = PP_ALIGN.CENTER

subtitle = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
p = subtitle.text_frame.paragraphs[0]
p.text = "Building the future, one solution at a time"
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0xCC, 0xCC, 0xCC)
p.alignment = PP_ALIGN.CENTER
'''
            }
        ]
    },

    "visual_emphasis": {
        "description": "Highlights, color boxes, accent elements",
        "when_to_use": "Call to action, key takeaways, important notices",
        "examples": [
            {
                "name": "accent_box",
                "code": '''
# Content with accent/highlight box
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(9), Inches(0.7))
p = title.text_frame.paragraphs[0]
p.text = "Key Takeaway"
p.font.size = Pt(32)
p.font.bold = True
p.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)

# Accent box (colored background)
accent_bg = slide.shapes.add_shape(1, Inches(0.5), Inches(1.5), Inches(9), Inches(3))
accent_bg.fill.solid()
accent_bg.fill.fore_color.rgb = RGBColor(0x00, 0x7A, 0xCC)
accent_bg.line.fill.background()

# Text inside accent box
accent_text = slide.shapes.add_textbox(Inches(0.8), Inches(2), Inches(8.4), Inches(2.5))
tf = accent_text.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = "Companies that invest in digital transformation see 2.5x higher revenue growth than their competitors."
p.font.size = Pt(24)
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
p.line_spacing = 1.4

# Source note below
source = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.5))
p = source.text_frame.paragraphs[0]
p.text = "Source: Industry Research 2024"
p.font.size = Pt(12)
p.font.color.rgb = RGBColor(0x7F, 0x8C, 0x8D)
'''
            },
            {
                "name": "top_accent_bar",
                "code": '''
# Slide with top accent bar
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(0x1A, 0x1A, 0x2E)

# Top accent bar
accent_bar = slide.shapes.add_shape(1, Inches(0), Inches(0), prs.slide_width, Inches(0.1))
accent_bar.fill.solid()
accent_bar.fill.fore_color.rgb = RGBColor(0x00, 0xD4, 0xAA)
accent_bar.line.fill.background()

# Title
title = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(0.8))
p = title.text_frame.paragraphs[0]
p.text = "Next Steps"
p.font.size = Pt(36)
p.font.bold = True
p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

# Content with numbered items
items = [
    "Schedule follow-up meeting",
    "Review proposal document",
    "Prepare technical requirements"
]

for i, item in enumerate(items):
    y = 1.8 + i * 1.2

    # Number circle
    circle = slide.shapes.add_shape(9, Inches(0.5), Inches(y), Inches(0.5), Inches(0.5))  # 9 = oval
    circle.fill.solid()
    circle.fill.fore_color.rgb = RGBColor(0x00, 0xD4, 0xAA)
    circle.line.fill.background()

    # Number text
    num = slide.shapes.add_textbox(Inches(0.5), Inches(y + 0.05), Inches(0.5), Inches(0.5))
    p = num.text_frame.paragraphs[0]
    p.text = str(i + 1)
    p.font.size = Pt(20)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x1A, 0x1A, 0x2E)
    from pptx.enum.text import PP_ALIGN
    p.alignment = PP_ALIGN.CENTER

    # Item text
    txt = slide.shapes.add_textbox(Inches(1.2), Inches(y + 0.05), Inches(8), Inches(0.5))
    p = txt.text_frame.paragraphs[0]
    p.text = item
    p.font.size = Pt(22)
    p.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
'''
            }
        ]
    }
}


def get_examples(category: str = None) -> dict:
    """Get slide code examples.

    Args:
        category: One of "text_layout", "number_highlight", "grid_layout",
                  "image_text", "visual_emphasis", or None for all

    Returns:
        Dict with examples for the category or all categories
    """
    if category and category in SLIDE_EXAMPLES:
        return {category: SLIDE_EXAMPLES[category]}
    return SLIDE_EXAMPLES


def get_all_categories() -> list:
    """Get list of available example categories."""
    return list(SLIDE_EXAMPLES.keys())
