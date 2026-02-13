"""
PowerPoint Presentation Tools

Tools for creating and editing PowerPoint presentations.
Uses Code Interpreter for python-pptx operations.
"""

import os
import logging
from typing import Dict, Any
from strands import tool, ToolContext
from skill import register_skill
from bedrock_agentcore.tools.code_interpreter_client import CodeInterpreter
from workspace import PowerPointManager

# Import utilities from lib
from .lib.ppt_utils import (
    validate_presentation_name,
    sanitize_presentation_name,
    get_code_interpreter_id,
    get_user_session_ids,
    save_ppt_artifact,
    get_file_compatibility_error,
    upload_ppt_helpers_to_ci,
    make_error_response,
)
from .lib.slide_examples import get_examples, get_all_categories

logger = logging.getLogger(__name__)


# Backward compatibility aliases
_validate_presentation_name = validate_presentation_name
_sanitize_presentation_name_for_bedrock = sanitize_presentation_name
_get_code_interpreter_id = get_code_interpreter_id
_get_user_session_ids = get_user_session_ids
_save_ppt_artifact = save_ppt_artifact
_get_file_compatibility_error_response = get_file_compatibility_error
_upload_ppt_helpers_to_ci = upload_ppt_helpers_to_ci


@tool
def get_slide_code_examples(category: str = "text_layout") -> Dict[str, Any]:
    """Get python-pptx code examples as reference for creating professional slides.

    These are reference examples with named color palettes and visual design patterns.
    Adapt them to your content needs. Also useful when debugging code errors.

    Args:
        category: Example category to retrieve:
            - "text_layout": Icon rows, accent bars, structured text (icon_text_rows, left_accent_bar)
            - "number_highlight": Hero stats, metric cards (hero_stat_dark, metric_cards)
            - "grid_layout": Comparison cards, feature columns (comparison_cards, three_column_feature)
            - "image_text": Half-bleed layouts, overlays (half_bleed_layout, full_bg_overlay)
            - "visual_emphasis": Quotes, process steps (quote_callout, process_steps)
            - "design_reference": Color palettes table and font pairing reference (text, not code)
            - "all": Get all categories

    Returns:
        Code examples with descriptions for the requested category
    """
    try:
        if category == "all":
            examples = get_examples()
        else:
            if category not in get_all_categories():
                return {
                    "content": [{
                        "text": f"Unknown category: {category}\n\n"
                               f"Available: {', '.join(get_all_categories())}"
                    }],
                    "status": "error"
                }
            examples = get_examples(category)

        # Format output
        output_parts = []
        for cat_name, cat_data in examples.items():
            output_parts.append(f"## {cat_name}\n")
            output_parts.append(f"**When to use:** {cat_data['when_to_use']}\n")

            for example in cat_data['examples']:
                output_parts.append(f"\n### {example['name']}\n")
                if 'text' in example:
                    # Text reference (design_reference category)
                    output_parts.append(f"{example['text'].strip()}\n")
                else:
                    # Code example
                    output_parts.append(f"```python\n{example['code'].strip()}\n```\n")

        return {
            "content": [{"text": "\n".join(output_parts)}],
            "status": "success",
            "metadata": {"category": category}
        }

    except Exception as e:
        logger.error(f"get_slide_code_examples error: {e}")
        return make_error_response(str(e))


@tool(context=True)
def list_my_powerpoint_presentations(tool_context: ToolContext) -> Dict[str, Any]:
    """List all PowerPoint presentations in workspace.

    Returns:
        Formatted list of presentations with metadata
    """
    try:
        logger.info("=== list_my_powerpoint_presentations called ===")

        # Get user and session IDs
        user_id, session_id = _get_user_session_ids(tool_context)

        # Initialize presentation manager
        ppt_manager = PowerPointManager(user_id, session_id)

        # List S3 documents
        documents = ppt_manager.list_s3_documents()

        # Format list
        workspace_list = ppt_manager.format_file_list(documents)

        return {
            "content": [{"text": workspace_list}],
            "status": "success",
            "metadata": {
                "count": len(documents),
                "presentations": [doc['filename'] for doc in documents]
            }
        }

    except Exception as e:
        logger.error(f"list_my_powerpoint_presentations error: {e}", exc_info=True)
        return {
            "content": [{
                "text": f"**Error listing presentations:** {str(e)}"
            }],
            "status": "error"
        }


@tool(context=True)
def get_presentation_layouts(
    presentation_name: str,
    tool_context: ToolContext
) -> Dict[str, Any]:
    """Get all available slide layouts from presentation.

    Returns complete list of layout names that can be used with add_slide.
    Use this before adding slides to get exact layout names.

    Args:
        presentation_name: Presentation name WITHOUT extension (e.g., "sales-deck")

    Example:
        get_presentation_layouts("sales-deck")
    """
    try:
        logger.info("=== get_presentation_layouts called ===")
        logger.info(f"Presentation: {presentation_name}")

        # Sanitize name for Bedrock API
        sanitized_name = _sanitize_presentation_name_for_bedrock(presentation_name)
        presentation_filename = f"{sanitized_name}.pptx"

        # Get user and session IDs
        user_id, session_id = _get_user_session_ids(tool_context)

        # Initialize presentation manager
        ppt_manager = PowerPointManager(user_id, session_id)

        # Load from S3
        try:
            pptx_bytes = ppt_manager.load_from_s3(presentation_filename)
        except FileNotFoundError:
            documents = ppt_manager.list_s3_documents()
            available = [doc['filename'] for doc in documents if doc['filename'].endswith('.pptx')]
            return {
                "content": [{
                    "text": f"**Presentation not found**: {presentation_filename}\n\n**Available presentations:**\n" + "\n".join([f"- {f}" for f in available]) if available else "No presentations found in workspace."
                }],
                "status": "error"
            }

        # Get Code Interpreter
        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {
                "content": [{
                    "text": "**Code Interpreter not configured**\n\nCODE_INTERPRETER_ID not found in environment or Parameter Store."
                }],
                "status": "error"
            }

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            # Upload presentation
            ppt_manager.upload_to_code_interpreter(code_interpreter, presentation_filename, pptx_bytes)

            # Generate code to list layouts
            list_layouts_code = f"""
from pptx import Presentation
import json

try:
    prs = Presentation('{presentation_filename}')

    layouts = []
    for idx, layout in enumerate(prs.slide_layouts):
        layouts.append({{
            'index': idx,
            'name': layout.name,
            'placeholder_count': len(layout.placeholders)
        }})

    result = {{
        'total_layouts': len(layouts),
        'layouts': layouts
    }}

    print(json.dumps(result, indent=2, ensure_ascii=False))

except Exception as e:
    raise ValueError(f"Failed to list layouts: {{str(e)}}")
""".strip()

            # Execute
            response = code_interpreter.invoke("executeCode", {
                "code": list_layouts_code,
                "language": "python",
                "clearContext": False
            })

            # Collect JSON output
            json_output = ""
            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"List layouts failed: {error_msg[:500]}")
                    code_interpreter.stop()
                    return {
                        "content": [{
                            "text": f"**Failed to list layouts**\n\n```\n{error_msg[:1000]}\n```"
                        }],
                        "status": "error"
                    }

                stdout = result.get("structuredContent", {}).get("stdout", "")
                if stdout:
                    json_output += stdout

            # Parse JSON result
            import json
            layout_data = json.loads(json_output)

            # Format output
            output_text = f"""📐 **Available Layouts**: {presentation_filename}

**Total layouts:** {layout_data['total_layouts']}

**Layout names (use exact names with add_slide):**
"""
            for layout in layout_data['layouts']:
                output_text += f"- \"{layout['name']}\"\n"

            code_interpreter.stop()
            return {
                "content": [{"text": output_text}],
                "status": "success",
                "metadata": {
                    "filename": presentation_filename,
                    "layouts": layout_data['layouts'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        finally:
            code_interpreter.stop()

    except Exception as e:
        logger.error(f"get_presentation_layouts error: {e}", exc_info=True)
        return {
            "content": [{
                "text": f"**Error getting layouts:** {str(e)}"
            }],
            "status": "error"
        }


@tool(context=True)
def analyze_presentation(
    presentation_name: str,
    tool_context: ToolContext,
    slide_index: int | None = None,
    include_notes: bool = False
) -> Dict[str, Any]:
    """Analyze presentation structure with element IDs and positions for editing.

    Returns layouts, slides, element IDs, text content, and spatial positions.
    Hidden slides are automatically excluded from analysis.

    Element Types:
    - text: Any text-editable shape (includes textbox, placeholder, autoshape, animation text, etc.)
    - picture: Image shapes
    - table: Table shapes (with cell data)
    - chart: Chart shapes (with data)
    - group: Grouped shapes
    - unknown: Other non-editable shapes

    Role Tags (for text elements):
    - [TITLE]: Title placeholder
    - [BODY]: Body/content placeholder
    - [FOOTER]: Footer placeholder
    - (no tag): Regular text element

    Args:
        presentation_name: Presentation name WITHOUT extension (e.g., "sales-deck")
        slide_index: Optional slide index (0-based). If provided, analyzes only that slide.
                     If None, analyzes entire presentation.
        include_notes: Whether to include speaker notes in output (default: False)

    Example:
        analyze_presentation("sales-deck")  # Analyze all visible slides
        analyze_presentation("sales-deck", slide_index=5)  # Analyze only slide 6
        analyze_presentation("sales-deck", include_notes=True)  # Include speaker notes

    Note: Element positions shown as (left", top") in inches.
    """
    try:
        logger.info("=== analyze_presentation called ===")
        logger.info(f"Presentation: {presentation_name}")

        # Sanitize name for Bedrock API
        sanitized_name = _sanitize_presentation_name_for_bedrock(presentation_name)
        presentation_filename = f"{sanitized_name}.pptx"

        # Get user and session IDs
        user_id, session_id = _get_user_session_ids(tool_context)

        # Initialize presentation manager
        ppt_manager = PowerPointManager(user_id, session_id)

        # Load from S3
        try:
            pptx_bytes = ppt_manager.load_from_s3(presentation_filename)
        except FileNotFoundError:
            documents = ppt_manager.list_s3_documents()
            available = [doc['filename'] for doc in documents if doc['filename'].endswith('.pptx')]
            return {
                "content": [{
                    "text": f"**Presentation not found**: {presentation_filename}\n\n**Available presentations:**\n" + "\n".join([f"- {f}" for f in available]) if available else "No presentations found in workspace."
                }],
                "status": "error"
            }

        # Get Code Interpreter
        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {
                "content": [{
                    "text": "**Code Interpreter not configured**\n\nCODE_INTERPRETER_ID not found in environment or Parameter Store."
                }],
                "status": "error"
            }

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            # Upload presentation
            ppt_manager.upload_to_code_interpreter(code_interpreter, presentation_filename, pptx_bytes)

            # Upload ppt_helpers (contains presentation_editor)
            _upload_ppt_helpers_to_ci(code_interpreter)

            # Generate analysis code
            from .lib.ppt_operations import generate_analyze_presentation_code
            analysis_code = generate_analyze_presentation_code(presentation_filename, slide_index)

            # Execute analysis
            response = code_interpreter.invoke("executeCode", {
                "code": analysis_code,
                "language": "python",
                "clearContext": False
            })

            # Collect JSON output
            json_output = ""
            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Analysis failed: {error_msg[:500]}")
                    code_interpreter.stop()

                    # Check if it's a file compatibility issue
                    if "AttributeError" in error_msg or "Cannot access" in error_msg:
                        return _get_file_compatibility_error_response(presentation_filename, error_msg, "analyze")

                    # Other errors
                    return {
                        "content": [{
                            "text": f"**Analysis failed**\n\n```\n{error_msg[:1000]}\n```"
                        }],
                        "status": "error"
                    }

                stdout = result.get("structuredContent", {}).get("stdout", "")
                if stdout:
                    json_output += stdout

            # Parse JSON result
            import json
            analysis = json.loads(json_output)

            # Format output text with detailed structure
            if slide_index is not None:
                # Single slide analysis
                output_text = f"""**Slide Analysis**: {presentation_filename} - Slide {slide_index + 1}

"""
            else:
                # Full presentation analysis
                output_text = f"""**Presentation Analysis**: {presentation_filename}

**Total slides:** {analysis['total_slides']}

**Slides:**
"""

            def format_ranges(indices):
                """Convert list of indices to range string: [0,1,2,5,6,10] -> '0-2, 5-6, 10'"""
                if not indices:
                    return ""

                indices = sorted(set(indices))
                ranges = []
                start = indices[0]
                end = indices[0]

                for i in range(1, len(indices)):
                    if indices[i] == end + 1:
                        end = indices[i]
                    else:
                        if start == end:
                            ranges.append(str(start))
                        else:
                            ranges.append(f"{start}-{end}")
                        start = indices[i]
                        end = indices[i]

                # Add last range
                if start == end:
                    ranges.append(str(start))
                else:
                    ranges.append(f"{start}-{end}")

                return ", ".join(ranges)

            for slide in analysis['slides']:
                # Slide header
                if slide_index is None:
                    output_text += f"\n**Slide {slide['index'] + 1}**\n"

                # Layout
                output_text += f"- **Layout:** {slide['layout']}\n"

                # Title
                if slide.get('title'):
                    output_text += f"- **Title:** {slide['title']}\n"

                # Show speaker notes if available and requested (full content)
                if include_notes and slide.get('notes'):
                    output_text += f"- **Notes:** \"{slide['notes']}\"\n"

                if 'elements' in slide:
                    editable_elements = []
                    excluded_indices = []

                    for elem in slide['elements']:
                        elem_type = elem['type']

                        # Exclude: non-editable types (picture, group, unknown)
                        if elem_type in ['picture', 'group', 'unknown']:
                            excluded_indices.append(elem['element_id'])
                            continue

                        # For text type: check if has paragraphs (even if empty)
                        if elem_type == 'text':
                            # Include if has paragraphs key (even if empty or all blank)
                            # Animation text objects often start empty but are still editable
                            if 'paragraphs' not in elem:
                                excluded_indices.append(elem['element_id'])
                                continue

                        # Keep: text (with text_frame), table, chart
                        editable_elements.append(elem)

                    # Display editable elements only (1 line per element)
                    if editable_elements:
                        output_text += "- **Elements:**\n"

                    for elem in editable_elements:
                        role_text = f" [{elem['role'].upper()}]" if elem.get('role') else ""
                        output_text += f"  - Element {elem['element_id']} ({elem['type']}{role_text})"

                        # Position (left, top only)
                        if 'position' in elem:
                            pos = elem['position']
                            output_text += f" @ ({pos['left']}\", {pos['top']}\")"

                        # Text element: show all paragraphs with format
                        if 'paragraphs' in elem:
                            if elem['paragraphs']:
                                # Collect all paragraph texts
                                all_texts = []
                                for para in elem['paragraphs']:
                                    para_text = para['text'].strip()
                                    if para_text:
                                        all_texts.append(para_text)

                                if all_texts:
                                    # Join with separator - show full text
                                    full_text = " | ".join(all_texts)
                                    output_text += f": \"{full_text}\""
                                else:
                                    output_text += ": (empty)"

                                # Format info from first paragraph (representative)
                                first_para = elem['paragraphs'][0]
                                if 'format' in first_para:
                                    fmt = first_para['format']
                                    format_parts = []
                                    if 'font_name' in fmt:
                                        format_parts.append(fmt['font_name'])
                                    if 'font_size' in fmt:
                                        format_parts.append(f"{fmt['font_size']}pt")
                                    if fmt.get('bold'):
                                        format_parts.append("Bold")
                                    if fmt.get('italic'):
                                        format_parts.append("Italic")
                                    if 'color' in fmt:
                                        format_parts.append(fmt['color'])

                                    if format_parts:
                                        output_text += f" [{', '.join(format_parts)}]"
                            else:
                                output_text += ": (empty)"

                        # Table: show size and first row
                        elif 'table_data' in elem:
                            output_text += f" [{elem['rows']}×{elem['cols']}]"
                            if elem['table_data']:
                                first_row = elem['table_data'][0]
                                row_preview = " | ".join([cell[:20] for cell in first_row[:3]])
                                if len(first_row) > 3:
                                    row_preview += "..."
                                output_text += f": \"{row_preview}\""

                        # Chart: show type and categories
                        elif 'chart_data' in elem:
                            chart = elem['chart_data']
                            chart_type = chart.get('chart_type', 'Unknown')
                            output_text += f" [{chart_type}]"
                            if 'categories' in chart:
                                cats_preview = ', '.join([str(c) for c in chart['categories'][:4]])
                                if len(chart['categories']) > 4:
                                    cats_preview += "..."
                                output_text += f": \"{cats_preview}\""

                        output_text += "\n"

                    # Summary of excluded elements
                    if excluded_indices:
                        output_text += f"  (Non-editable: {format_ranges(excluded_indices)})\n"

            code_interpreter.stop()
            return {
                "content": [{"text": output_text}],
                "status": "success",
                "metadata": {
                    "filename": presentation_filename,
                    "slide_index": slide_index,
                    "analysis": analysis,
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        finally:
            code_interpreter.stop()

    except Exception as e:
        logger.error(f"analyze_presentation error: {e}", exc_info=True)
        return {
            "content": [{
                "text": f"**Error analyzing presentation:** {str(e)}"
            }],
            "status": "error"
        }


@tool(context=True)
def update_slide_content(
    presentation_name: str,
    slide_updates: list,
    output_name: str,
    tool_context: ToolContext
) -> Dict[str, Any]:
    """Update one or more slides with operations in a single call.

    Args:
        presentation_name: Source presentation name WITHOUT extension
        slide_updates: List of slide update dicts with structure:
            [
                {
                    "slide_index": int (0-based),
                    "operations": [
                        {"action": "set_text", "element_id": int, "text": str},
                        {"action": "replace_text", "element_id": int, "find": str, "replace": str},
                        {"action": "replace_image", "element_id": int, "image_path": str}
                    ]
                }
            ]
        output_name: Output presentation name WITHOUT extension (must differ from source)

    Supported Actions:
        - set_text: Replace entire element text (works on type: text)
        - replace_text: Find and replace within element (works on type: text)
        - replace_image: Replace element image (works on type: picture)

    Element IDs:
        - Get element_id from analyze_presentation output
        - Element IDs are 0-based indices within each slide
        - text type: Use set_text or replace_text
        - picture type: Use replace_image

    Example (single slide):
        update_slide_content("sales-deck",
            slide_updates=[{
                "slide_index": 5,
                "operations": [
                    {"action": "set_text", "element_id": 0, "text": "Q4 Results"},
                    {"action": "replace_text", "element_id": 1, "find": "2024", "replace": "2025"}
                ]
            }],
            output_name="sales-deck-v2")

    Example (multiple slides):
        update_slide_content("sales-deck",
            slide_updates=[
                {"slide_index": 2, "operations": [{"action": "set_text", "element_id": 0, "text": "New Title"}]},
                {"slide_index": 5, "operations": [
                    {"action": "replace_text", "element_id": 1, "find": "Quality assessment", "replace": "Quality evaluation"},
                    {"action": "replace_image", "element_id": 3, "image_path": "new-chart.png"}
                ]},
                {"slide_index": 8, "operations": [{"action": "set_text", "element_id": 1, "text": "Point 1\\nPoint 2\\nPoint 3"}]}
            ],
            output_name="sales-deck-v2")

    Notes:
        - Multi-line text (\\n) creates multiple paragraphs automatically
        - Formatting is preserved from original text
        - Use replace_text for long texts where you only need to change specific portions
        - Batch all modifications into ONE call to avoid data loss from parallel execution
    """
    try:
        logger.info("=== update_slide_content called ===")
        logger.info(f"Source: {presentation_name}, Updates: {len(slide_updates)} slide(s)")

        # Validate names
        is_valid_source, error_msg_source = _validate_presentation_name(presentation_name)
        if not is_valid_source:
            return {
                "content": [{
                    "text": f"**Invalid source name**: {presentation_name}\n\n{error_msg_source}"
                }],
                "status": "error"
            }

        is_valid_output, error_msg_output = _validate_presentation_name(output_name)
        if not is_valid_output:
            return {
                "content": [{
                    "text": f"**Invalid output name**: {output_name}\n\n{error_msg_output}"
                }],
                "status": "error"
            }

        if presentation_name == output_name:
            return {
                "content": [{
                    "text": f"**Output name must be different from source name**\n\nSource: {presentation_name}\nOutput: {output_name}\n\nThis preserves the original file."
                }],
                "status": "error"
            }

        # Validate slide_updates
        if not slide_updates or not isinstance(slide_updates, list):
            return {
                "content": [{
                    "text": "**Invalid slide_updates**: Must provide a non-empty list of slide update dicts"
                }],
                "status": "error"
            }

        # Validate each slide update
        for idx, update in enumerate(slide_updates):
            if not isinstance(update, dict):
                return {
                    "content": [{
                        "text": f"**Invalid slide_updates[{idx}]**: Must be a dict with 'slide_index' and 'operations'"
                    }],
                    "status": "error"
                }

            if 'slide_index' not in update or 'operations' not in update:
                return {
                    "content": [{
                        "text": f"**Invalid slide_updates[{idx}]**: Must have 'slide_index' and 'operations' keys"
                    }],
                    "status": "error"
                }

            if not isinstance(update['operations'], list) or not update['operations']:
                return {
                    "content": [{
                        "text": f"**Invalid slide_updates[{idx}]['operations']**: Must be a non-empty list"
                    }],
                    "status": "error"
                }

        # Add extensions
        source_filename = f"{presentation_name}.pptx"
        output_filename = f"{output_name}.pptx"

        # Get user and session IDs
        user_id, session_id = _get_user_session_ids(tool_context)

        # Initialize presentation manager
        ppt_manager = PowerPointManager(user_id, session_id)

        # Load source from S3
        try:
            source_bytes = ppt_manager.load_from_s3(source_filename)
        except FileNotFoundError:
            documents = ppt_manager.list_s3_documents()
            available = [doc['filename'] for doc in documents if doc['filename'].endswith('.pptx')]
            return {
                "content": [{
                    "text": f"**Source presentation not found**: {source_filename}\n\n**Available:**\n" + "\n".join([f"- {f}" for f in available])
                }],
                "status": "error"
            }

        # Get Code Interpreter
        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {
                "content": [{
                    "text": "**Code Interpreter not configured**"
                }],
                "status": "error"
            }

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            # Upload source
            ppt_manager.upload_to_code_interpreter(code_interpreter, source_filename, source_bytes)

            # Upload workspace images
            loaded_images = ppt_manager.load_workspace_images_to_ci(code_interpreter)

            # Upload ppt_helpers (contains PresentationEditor)
            _upload_ppt_helpers_to_ci(code_interpreter)

            # Generate safe code from slide updates
            from .lib.ppt_operations import generate_batch_update_slides_code
            safe_code = generate_batch_update_slides_code(
                source_filename,
                output_filename,
                slide_updates
            )

            logger.info(f"Generated code:\n{safe_code[:500]}...")

            # Execute
            response = code_interpreter.invoke("executeCode", {
                "code": safe_code,
                "language": "python",
                "clearContext": False
            })

            # Check for errors
            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Update failed: {error_msg[:500]}")
                    code_interpreter.stop()

                    # Check if it's a file compatibility issue
                    if "AttributeError" in error_msg or "Cannot access" in error_msg:
                        return _get_file_compatibility_error_response(source_filename, error_msg, "update")

                    # Other errors
                    return {
                        "content": [{
                            "text": f"**Slide update failed**\n\n**Error:**\n```\n{error_msg[:1000]}\n```"
                        }],
                        "status": "error"
                    }

            # Download result
            output_ci_path = ppt_manager.get_ci_path(output_filename)
            file_bytes = ppt_manager.download_from_code_interpreter(code_interpreter, output_ci_path)

            if not file_bytes:
                code_interpreter.stop()
                return {
                    "content": [{
                        "text": "**Failed to retrieve updated presentation**"
                    }],
                    "status": "error"
                }

            # Save to S3
            s3_info = ppt_manager.save_to_s3(output_filename, file_bytes)

            # Save as artifact for Canvas display
            _save_ppt_artifact(
                tool_context=tool_context,
                filename=output_filename,
                s3_url=s3_info['s3_url'],
                size_kb=s3_info['size_kb'],
                tool_name='update_slide_content',
                user_id=user_id,
                session_id=session_id
            )

            # Get updated workspace list
            documents = ppt_manager.list_s3_documents()
            other_files_count = len([d for d in documents if d['filename'] != output_filename])

            # Calculate total operations
            total_operations = sum(len(update['operations']) for update in slide_updates)
            slide_numbers = [update['slide_index'] + 1 for update in slide_updates]

            # Success message
            if len(slide_updates) == 1:
                success_msg = f"""**Slide updated successfully!**

**Original:** {source_filename} (preserved)
**Updated:** {output_filename}
**Slide:** {slide_numbers[0]}
**Operations:** {total_operations}
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}

**Next steps:**
- Use `update_slide_content` again for more edits
- Use `analyze_presentation` to verify changes
"""
            else:
                slides_str = ", ".join(str(n) for n in sorted(slide_numbers))
                success_msg = f"""**Slides updated successfully!**

**Original:** {source_filename} (preserved)
**Updated:** {output_filename}
**Slides:** {slides_str}
**Total operations:** {total_operations}
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}

**Next steps:**
- Use `update_slide_content` again for more edits
- Use `analyze_presentation` to verify changes
"""

            code_interpreter.stop()
            return {
                "content": [{"text": success_msg}],
                "status": "success",
                "metadata": {
                    "source_filename": source_filename,
                    "filename": output_filename,
                    "slide_updates": slide_updates,
                    "total_operations": total_operations,
                    "size_kb": s3_info['size_kb'],
                    "s3_key": s3_info['s3_key'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        except Exception as e:
            code_interpreter.stop()
            logger.error(f"Slide update error: {e}", exc_info=True)
            return {
                "content": [{
                    "text": f"**Slide update failed**\n\n**Error:** {str(e)}"
                }],
                "status": "error"
            }

    except Exception as e:
        logger.error(f"update_slide_content error: {e}", exc_info=True)
        return {
            "content": [{
                "text": f"**Error:** {str(e)}"
            }],
            "status": "error"
        }


@tool(context=True)
def add_slide(
    presentation_name: str,
    layout_name: str,
    position: int,
    output_name: str,
    tool_context: ToolContext,
    custom_code: str | None = None
) -> Dict[str, Any]:
    """Add new slide with optional custom python-pptx code.

    Use get_slide_code_examples() for custom_code reference.

    Args:
        presentation_name: Source presentation name without extension
        layout_name: Layout name (use get_presentation_layouts() to get exact names)
        position: Insert position (0-based, -1 to append)
        output_name: Output name without extension (must differ from source)
        custom_code: Optional python-pptx code. Available: slide, prs, Inches, Pt, RGBColor

    Example:
        add_slide("deck", "Blank", position=0, output_name="deck-v2",
                  custom_code="slide.shapes.title.text = 'New Slide'")
    """
    try:
        logger.info("=== add_slide called ===")
        logger.info(f"Source: {presentation_name}, Layout: {layout_name}, Position: {position}")

        # Validate names
        is_valid_source, error_msg_source = _validate_presentation_name(presentation_name)
        if not is_valid_source:
            return {"content": [{"text": f"**Invalid source name**: {presentation_name}\n\n{error_msg_source}"}], "status": "error"}

        is_valid_output, error_msg_output = _validate_presentation_name(output_name)
        if not is_valid_output:
            return {"content": [{"text": f"**Invalid output name**: {output_name}\n\n{error_msg_output}"}], "status": "error"}

        if presentation_name == output_name:
            return {"content": [{"text": "**Output name must be different from source name**"}], "status": "error"}

        source_filename = f"{presentation_name}.pptx"
        output_filename = f"{output_name}.pptx"

        user_id, session_id = _get_user_session_ids(tool_context)
        ppt_manager = PowerPointManager(user_id, session_id)

        try:
            source_bytes = ppt_manager.load_from_s3(source_filename)
        except FileNotFoundError:
            return {"content": [{"text": f"**Source presentation not found**: {source_filename}"}], "status": "error"}

        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {"content": [{"text": "**Code Interpreter not configured**"}], "status": "error"}

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            ppt_manager.upload_to_code_interpreter(code_interpreter, source_filename, source_bytes)

            # Load workspace images if custom_code provided (might use images)
            if custom_code:
                ppt_manager.load_workspace_images_to_ci(code_interpreter)

            from .lib.ppt_operations import generate_add_slide_code
            safe_code = generate_add_slide_code(source_filename, output_filename, layout_name, position, custom_code)

            response = code_interpreter.invoke("executeCode", {"code": safe_code, "language": "python", "clearContext": False})

            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Add slide failed: {error_msg[:500]}")
                    code_interpreter.stop()

                    # Check if it's a file compatibility issue
                    if "AttributeError" in error_msg or "Cannot access" in error_msg:
                        return _get_file_compatibility_error_response(source_filename, error_msg, "add slide to")

                    # Other errors
                    return {"content": [{"text": f"**Failed to add slide**\n\n```\n{error_msg[:1000]}\n```"}], "status": "error"}

            output_ci_path = ppt_manager.get_ci_path(output_filename)
            file_bytes = ppt_manager.download_from_code_interpreter(code_interpreter, output_ci_path)

            if not file_bytes:
                code_interpreter.stop()
                return {"content": [{"text": "**Failed to retrieve presentation**"}], "status": "error"}

            s3_info = ppt_manager.save_to_s3(output_filename, file_bytes)

            # Save as artifact for Canvas display
            _save_ppt_artifact(
                tool_context=tool_context,
                filename=output_filename,
                s3_url=s3_info['s3_url'],
                size_kb=s3_info['size_kb'],
                tool_name='add_slide',
                user_id=user_id,
                session_id=session_id
            )

            documents = ppt_manager.list_s3_documents()
            other_files_count = len([d for d in documents if d['filename'] != output_filename])

            success_msg = f"""**Slide added successfully!**

**Original:** {source_filename} (preserved)
**Updated:** {output_filename}
**Layout:** {layout_name}
**Position:** {position if position >= 0 else 'end'}
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}"""

            code_interpreter.stop()
            return {
                "content": [{"text": success_msg}],
                "status": "success",
                "metadata": {
                    "filename": output_filename,
                    "layout_name": layout_name,
                    "position": position,
                    "size_kb": s3_info['size_kb'],
                    "s3_key": s3_info['s3_key'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        except Exception as e:
            code_interpreter.stop()
            raise e

    except Exception as e:
        logger.error(f"add_slide error: {e}", exc_info=True)
        return {"content": [{"text": f"**Error:** {str(e)}"}], "status": "error"}


@tool(context=True)
def delete_slides(
    presentation_name: str,
    slide_indices: list,
    output_name: str,
    tool_context: ToolContext
) -> Dict[str, Any]:
    """Delete slides by indices.

    Args:
        presentation_name: Source presentation name WITHOUT extension
        slide_indices: List of slide indices to delete (0-based, e.g., [2, 5, 10])
        output_name: Output presentation name WITHOUT extension (must differ from source)

    Example:
        delete_slides("sales-deck", slide_indices=[5, 12], output_name="sales-deck-cleaned")
    """
    try:
        logger.info("=== delete_slides called ===")
        logger.info(f"Source: {presentation_name}, Indices: {slide_indices}")

        is_valid_source, error_msg_source = _validate_presentation_name(presentation_name)
        if not is_valid_source:
            return {"content": [{"text": f"**Invalid source name**: {presentation_name}\n\n{error_msg_source}"}], "status": "error"}

        is_valid_output, error_msg_output = _validate_presentation_name(output_name)
        if not is_valid_output:
            return {"content": [{"text": f"**Invalid output name**: {output_name}\n\n{error_msg_output}"}], "status": "error"}

        if presentation_name == output_name:
            return {"content": [{"text": "**Output name must be different from source name**"}], "status": "error"}

        if not slide_indices or not isinstance(slide_indices, list):
            return {"content": [{"text": "**Invalid slide_indices**: Must provide a non-empty list"}], "status": "error"}

        source_filename = f"{presentation_name}.pptx"
        output_filename = f"{output_name}.pptx"

        user_id, session_id = _get_user_session_ids(tool_context)
        ppt_manager = PowerPointManager(user_id, session_id)

        try:
            source_bytes = ppt_manager.load_from_s3(source_filename)
        except FileNotFoundError:
            return {"content": [{"text": f"**Source presentation not found**: {source_filename}"}], "status": "error"}

        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {"content": [{"text": "**Code Interpreter not configured**"}], "status": "error"}

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            ppt_manager.upload_to_code_interpreter(code_interpreter, source_filename, source_bytes)

            from .lib.ppt_operations import generate_delete_slides_code
            safe_code = generate_delete_slides_code(source_filename, output_filename, slide_indices)

            response = code_interpreter.invoke("executeCode", {"code": safe_code, "language": "python", "clearContext": False})

            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Delete slides failed: {error_msg[:500]}")
                    code_interpreter.stop()

                    # Check if it's a file compatibility issue
                    if "AttributeError" in error_msg or "Cannot access" in error_msg:
                        return _get_file_compatibility_error_response(source_filename, error_msg, "delete slides from")

                    # Other errors
                    return {"content": [{"text": f"**Failed to delete slides**\n\n```\n{error_msg[:1000]}\n```"}], "status": "error"}

            output_ci_path = ppt_manager.get_ci_path(output_filename)
            file_bytes = ppt_manager.download_from_code_interpreter(code_interpreter, output_ci_path)

            if not file_bytes:
                code_interpreter.stop()
                return {"content": [{"text": "**Failed to retrieve presentation**"}], "status": "error"}

            s3_info = ppt_manager.save_to_s3(output_filename, file_bytes)

            # Save as artifact for Canvas display
            _save_ppt_artifact(
                tool_context=tool_context,
                filename=output_filename,
                s3_url=s3_info['s3_url'],
                size_kb=s3_info['size_kb'],
                tool_name='delete_slides',
                user_id=user_id,
                session_id=session_id
            )

            documents = ppt_manager.list_s3_documents()
            other_files_count = len([d for d in documents if d['filename'] != output_filename])

            success_msg = f"""**Slides deleted successfully!**

**Original:** {source_filename} (preserved)
**Updated:** {output_filename}
**Deleted slides:** {', '.join([str(i+1) for i in sorted(slide_indices)])}
**Count:** {len(slide_indices)} slide(s)
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}"""

            code_interpreter.stop()
            return {
                "content": [{"text": success_msg}],
                "status": "success",
                "metadata": {
                    "filename": output_filename,
                    "deleted_count": len(slide_indices),
                    "size_kb": s3_info['size_kb'],
                    "s3_key": s3_info['s3_key'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        except Exception as e:
            code_interpreter.stop()
            raise e

    except Exception as e:
        logger.error(f"delete_slides error: {e}", exc_info=True)
        return {"content": [{"text": f"**Error:** {str(e)}"}], "status": "error"}


@tool(context=True)
def move_slide(
    presentation_name: str,
    from_index: int,
    to_index: int,
    output_name: str,
    tool_context: ToolContext
) -> Dict[str, Any]:
    """Move slide from one position to another.

    Args:
        presentation_name: Source presentation name WITHOUT extension
        from_index: Source position (0-based)
        to_index: Target position (0-based)
        output_name: Output presentation name WITHOUT extension (must differ from source)

    Example:
        move_slide("sales-deck", from_index=10, to_index=2, output_name="sales-deck-reordered")
    """
    try:
        logger.info("=== move_slide called ===")
        logger.info(f"Source: {presentation_name}, From: {from_index}, To: {to_index}")

        is_valid_source, error_msg_source = _validate_presentation_name(presentation_name)
        if not is_valid_source:
            return {"content": [{"text": f"**Invalid source name**: {presentation_name}\n\n{error_msg_source}"}], "status": "error"}

        is_valid_output, error_msg_output = _validate_presentation_name(output_name)
        if not is_valid_output:
            return {"content": [{"text": f"**Invalid output name**: {output_name}\n\n{error_msg_output}"}], "status": "error"}

        if presentation_name == output_name:
            return {"content": [{"text": "**Output name must be different from source name**"}], "status": "error"}

        source_filename = f"{presentation_name}.pptx"
        output_filename = f"{output_name}.pptx"

        user_id, session_id = _get_user_session_ids(tool_context)
        ppt_manager = PowerPointManager(user_id, session_id)

        try:
            source_bytes = ppt_manager.load_from_s3(source_filename)
        except FileNotFoundError:
            return {"content": [{"text": f"**Source presentation not found**: {source_filename}"}], "status": "error"}

        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {"content": [{"text": "**Code Interpreter not configured**"}], "status": "error"}

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            ppt_manager.upload_to_code_interpreter(code_interpreter, source_filename, source_bytes)

            from .lib.ppt_operations import generate_move_slide_code
            safe_code = generate_move_slide_code(source_filename, output_filename, from_index, to_index)

            response = code_interpreter.invoke("executeCode", {"code": safe_code, "language": "python", "clearContext": False})

            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Move slide failed: {error_msg[:500]}")
                    code_interpreter.stop()

                    # Check if it's a file compatibility issue
                    if "AttributeError" in error_msg or "Cannot access" in error_msg:
                        return _get_file_compatibility_error_response(source_filename, error_msg, "move slide in")

                    # Other errors
                    return {"content": [{"text": f"**Failed to move slide**\n\n```\n{error_msg[:1000]}\n```"}], "status": "error"}

            output_ci_path = ppt_manager.get_ci_path(output_filename)
            file_bytes = ppt_manager.download_from_code_interpreter(code_interpreter, output_ci_path)

            if not file_bytes:
                code_interpreter.stop()
                return {"content": [{"text": "**Failed to retrieve presentation**"}], "status": "error"}

            s3_info = ppt_manager.save_to_s3(output_filename, file_bytes)

            # Save as artifact for Canvas display
            _save_ppt_artifact(
                tool_context=tool_context,
                filename=output_filename,
                s3_url=s3_info['s3_url'],
                size_kb=s3_info['size_kb'],
                tool_name='move_slide',
                user_id=user_id,
                session_id=session_id
            )

            documents = ppt_manager.list_s3_documents()
            other_files_count = len([d for d in documents if d['filename'] != output_filename])

            success_msg = f"""**Slide moved successfully!**

**Original:** {source_filename} (preserved)
**Updated:** {output_filename}
**Moved:** Slide {from_index + 1} → Position {to_index + 1}
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}"""

            code_interpreter.stop()
            return {
                "content": [{"text": success_msg}],
                "status": "success",
                "metadata": {
                    "filename": output_filename,
                    "from_index": from_index,
                    "to_index": to_index,
                    "size_kb": s3_info['size_kb'],
                    "s3_key": s3_info['s3_key'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        except Exception as e:
            code_interpreter.stop()
            raise e

    except Exception as e:
        logger.error(f"move_slide error: {e}", exc_info=True)
        return {"content": [{"text": f"**Error:** {str(e)}"}], "status": "error"}


@tool(context=True)
def duplicate_slide(
    presentation_name: str,
    slide_index: int,
    position: int,
    output_name: str,
    tool_context: ToolContext
) -> Dict[str, Any]:
    """Duplicate slide to specified position.

    Args:
        presentation_name: Source presentation name WITHOUT extension
        slide_index: Index of slide to duplicate (0-based)
        position: Position for duplicate (0-based, -1 for after original)
        output_name: Output presentation name WITHOUT extension (must differ from source)

    Example:
        duplicate_slide("sales-deck", slide_index=5, position=-1, output_name="sales-deck-v2")
    """
    try:
        logger.info("=== duplicate_slide called ===")
        logger.info(f"Source: {presentation_name}, Slide: {slide_index}, Position: {position}")

        is_valid_source, error_msg_source = _validate_presentation_name(presentation_name)
        if not is_valid_source:
            return {"content": [{"text": f"**Invalid source name**: {presentation_name}\n\n{error_msg_source}"}], "status": "error"}

        is_valid_output, error_msg_output = _validate_presentation_name(output_name)
        if not is_valid_output:
            return {"content": [{"text": f"**Invalid output name**: {output_name}\n\n{error_msg_output}"}], "status": "error"}

        if presentation_name == output_name:
            return {"content": [{"text": "**Output name must be different from source name**"}], "status": "error"}

        source_filename = f"{presentation_name}.pptx"
        output_filename = f"{output_name}.pptx"

        user_id, session_id = _get_user_session_ids(tool_context)
        ppt_manager = PowerPointManager(user_id, session_id)

        try:
            source_bytes = ppt_manager.load_from_s3(source_filename)
        except FileNotFoundError:
            return {"content": [{"text": f"**Source presentation not found**: {source_filename}"}], "status": "error"}

        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {"content": [{"text": "**Code Interpreter not configured**"}], "status": "error"}

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            ppt_manager.upload_to_code_interpreter(code_interpreter, source_filename, source_bytes)

            from .lib.ppt_operations import generate_duplicate_slide_code
            safe_code = generate_duplicate_slide_code(source_filename, output_filename, slide_index, position)

            response = code_interpreter.invoke("executeCode", {"code": safe_code, "language": "python", "clearContext": False})

            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Duplicate slide failed: {error_msg[:500]}")
                    code_interpreter.stop()

                    # Check if it's a file compatibility issue
                    if "AttributeError" in error_msg or "Cannot access" in error_msg:
                        return _get_file_compatibility_error_response(source_filename, error_msg, "duplicate slide in")

                    # Other errors
                    return {"content": [{"text": f"**Failed to duplicate slide**\n\n```\n{error_msg[:1000]}\n```"}], "status": "error"}

            output_ci_path = ppt_manager.get_ci_path(output_filename)
            file_bytes = ppt_manager.download_from_code_interpreter(code_interpreter, output_ci_path)

            if not file_bytes:
                code_interpreter.stop()
                return {"content": [{"text": "**Failed to retrieve presentation**"}], "status": "error"}

            s3_info = ppt_manager.save_to_s3(output_filename, file_bytes)

            # Save as artifact for Canvas display
            _save_ppt_artifact(
                tool_context=tool_context,
                filename=output_filename,
                s3_url=s3_info['s3_url'],
                size_kb=s3_info['size_kb'],
                tool_name='duplicate_slide',
                user_id=user_id,
                session_id=session_id
            )

            documents = ppt_manager.list_s3_documents()
            other_files_count = len([d for d in documents if d['filename'] != output_filename])

            success_msg = f"""**Slide duplicated successfully!**

**Original:** {source_filename} (preserved)
**Updated:** {output_filename}
**Duplicated:** Slide {slide_index + 1}
**New position:** {position if position >= 0 else f'{slide_index + 2} (after original)'}
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}"""

            code_interpreter.stop()
            return {
                "content": [{"text": success_msg}],
                "status": "success",
                "metadata": {
                    "filename": output_filename,
                    "slide_index": slide_index,
                    "position": position,
                    "size_kb": s3_info['size_kb'],
                    "s3_key": s3_info['s3_key'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        except Exception as e:
            code_interpreter.stop()
            raise e

    except Exception as e:
        logger.error(f"duplicate_slide error: {e}", exc_info=True)
        return {
            "content": [{
                "text": f"**Error duplicating slide:** {str(e)}"
            }],
            "status": "error"
        }


@tool(context=True)
def update_slide_notes(
    presentation_name: str,
    slide_index: int,
    notes_text: str,
    output_name: str,
    tool_context: ToolContext
) -> Dict[str, Any]:
    """Update speaker notes for a specific slide.

    Replaces entire notes content for the specified slide.
    Use when user wants to review or modify presentation notes.

    Args:
        presentation_name: Source presentation name WITHOUT extension
        slide_index: Slide index (0-based)
        notes_text: New notes content (use \n for multi-line)
        output_name: Output presentation name WITHOUT extension (must differ from source)

    Example:
        update_slide_notes("sales-deck", slide_index=5,
                          notes_text="Remember to emphasize ROI\\nMention competitor comparison",
                          output_name="sales-deck-v2")
    """
    try:
        logger.info("=== update_slide_notes called ===")
        logger.info(f"Source: {presentation_name}, Slide: {slide_index}")

        is_valid_source, error_msg_source = _validate_presentation_name(presentation_name)
        if not is_valid_source:
            return {"content": [{"text": f"**Invalid source name**: {presentation_name}\n\n{error_msg_source}"}], "status": "error"}

        is_valid_output, error_msg_output = _validate_presentation_name(output_name)
        if not is_valid_output:
            return {"content": [{"text": f"**Invalid output name**: {output_name}\n\n{error_msg_output}"}], "status": "error"}

        if presentation_name == output_name:
            return {"content": [{"text": "**Output name must be different from source name**"}], "status": "error"}

        source_filename = f"{presentation_name}.pptx"
        output_filename = f"{output_name}.pptx"

        user_id, session_id = _get_user_session_ids(tool_context)
        ppt_manager = PowerPointManager(user_id, session_id)

        try:
            source_bytes = ppt_manager.load_from_s3(source_filename)
        except FileNotFoundError:
            return {"content": [{"text": f"**Source presentation not found**: {source_filename}"}], "status": "error"}

        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {"content": [{"text": "**Code Interpreter not configured**"}], "status": "error"}

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(code_interpreter_id)

        try:
            ci_path = ppt_manager.upload_to_code_interpreter(code_interpreter, source_filename, source_bytes)

            from .lib.ppt_operations import generate_update_notes_code
            safe_code = generate_update_notes_code(source_filename, output_filename, slide_index, notes_text)

            response = code_interpreter.invoke("executeCode", {"code": safe_code, "language": "python", "clearContext": False})

            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Update notes failed: {error_msg[:500]}")
                    code_interpreter.stop()

                    if "AttributeError" in error_msg or "Cannot access" in error_msg:
                        return _get_file_compatibility_error_response(source_filename, error_msg, "update notes in")

                    return {"content": [{"text": f"**Failed to update notes**\n\n```\n{error_msg[:1000]}\n```"}], "status": "error"}

            output_ci_path = ppt_manager.get_ci_path(output_filename)
            file_bytes = ppt_manager.download_from_code_interpreter(code_interpreter, output_ci_path)

            if not file_bytes:
                code_interpreter.stop()
                return {"content": [{"text": "**Failed to retrieve presentation**"}], "status": "error"}

            s3_info = ppt_manager.save_to_s3(output_filename, file_bytes)

            # Save as artifact for Canvas display
            _save_ppt_artifact(
                tool_context=tool_context,
                filename=output_filename,
                s3_url=s3_info['s3_url'],
                size_kb=s3_info['size_kb'],
                tool_name='update_slide_notes',
                user_id=user_id,
                session_id=session_id
            )

            documents = ppt_manager.list_s3_documents()
            other_files_count = len([d for d in documents if d['filename'] != output_filename])

            notes_preview = notes_text[:100] + "..." if len(notes_text) > 100 else notes_text
            success_msg = f"""**Slide notes updated successfully!**

**Original:** {source_filename} (preserved)
**Updated:** {output_filename}
**Slide:** {slide_index + 1}
**Notes:** "{notes_preview}"
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}"""

            code_interpreter.stop()
            return {
                "content": [{"text": success_msg}],
                "status": "success",
                "metadata": {
                    "filename": output_filename,
                    "slide_index": slide_index,
                    "size_kb": s3_info['size_kb'],
                    "s3_key": s3_info['s3_key'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        except Exception as e:
            code_interpreter.stop()
            raise e

    except Exception as e:
        logger.error(f"update_slide_notes error: {e}", exc_info=True)
        return {
            "content": [{
                "text": f"**Error updating slide notes:** {str(e)}"
            }],
            "status": "error"
        }


@tool(context=True)
def create_presentation(
    presentation_name: str,
    slides: list | str | None,
    tool_context: ToolContext,
    template_name: str | None = None,
) -> Dict[str, Any]:
    """Create presentation with custom-designed slides (16:9 widescreen).

    Use get_slide_code_examples() to get reference code for different slide types.

    Args:
        presentation_name: Output name without extension (e.g., "sales-deck")
        slides: List of slide definitions with custom_code, or None for blank
        template_name: Optional template presentation name (uses template's aspect ratio)

    Format:
        slides = [{"custom_code": "...python-pptx code..."}]

    Slide size:
        16:9 widescreen (13.333" x 7.5") - standard for modern presentations.
        If template_name is provided, the template's aspect ratio is used instead.

    Available in custom_code:
        prs, slide, slide_width, slide_height, Inches, Pt, RGBColor, PP_ALIGN, MSO_SHAPE
    """
    try:
        logger.info("=== create_presentation called ===")
        logger.info(f"Name: {presentation_name}, Has slides: {slides is not None}")

        # Parse slides if it's a JSON string (LLM sometimes sends JSON as string)
        if isinstance(slides, str):
            import json
            import re

            logger.info(f"Received slides as string, length: {len(slides)}")
            logger.debug(f"Slides string (first 500 chars): {slides[:500]}")

            try:
                # Try standard JSON parsing first
                slides = json.loads(slides)
                logger.info("Successfully parsed slides from JSON string")
            except json.JSONDecodeError as e:
                # Try fixing common JSON issues
                logger.warning(f"Initial JSON parse failed: {str(e)}, attempting fixes...")

                try:
                    # Remove trailing commas before ] or }
                    fixed_json = re.sub(r',(\s*[}\]])', r'\1', slides)

                    # Remove comments (// style)
                    fixed_json = re.sub(r'//.*?$', '', fixed_json, flags=re.MULTILINE)

                    # Try parsing again
                    slides = json.loads(fixed_json)
                    logger.info("Successfully parsed slides after fixing JSON issues")
                except json.JSONDecodeError as e2:
                    # Log the problematic JSON for debugging
                    logger.error(f"JSON parse failed even after fixes. Error: {str(e2)}")
                    logger.error(f"Problematic JSON snippet around error position: {slides[max(0, e2.pos-50):min(len(slides), e2.pos+50)]}")

                    return {
                        "content": [{
                            "text": f"**Invalid JSON format for slides**\n\nError: {str(e2)}\n\n**Position**: Line {e2.lineno}, Column {e2.colno}\n\n**Hint**: Check for trailing commas, unescaped quotes, or invalid characters around the error position.\n\nPlease provide a valid JSON array."
                        }],
                        "status": "error"
                    }

        # Validate name
        is_valid, error_msg = _validate_presentation_name(presentation_name)
        if not is_valid:
            return {"content": [{"text": f"**Invalid presentation name**: {presentation_name}\n\n{error_msg}"}], "status": "error"}

        presentation_filename = f"{presentation_name}.pptx"

        user_id, session_id = _get_user_session_ids(tool_context)
        ppt_manager = PowerPointManager(user_id, session_id)

        # Check if file already exists
        try:
            ppt_manager.load_from_s3(presentation_filename)
            return {
                "content": [{
                    "text": f"**Presentation already exists**: {presentation_filename}\n\nPlease use a different name or delete the existing file first."
                }],
                "status": "error"
            }
        except FileNotFoundError:
            pass  # Good, file doesn't exist

        code_interpreter_id = _get_code_interpreter_id()
        if not code_interpreter_id:
            return {"content": [{"text": "**Code Interpreter not configured**"}], "status": "error"}

        region = os.getenv('AWS_REGION', 'us-west-2')
        code_interpreter = CodeInterpreter(region)
        code_interpreter.start(identifier=code_interpreter_id)

        try:
            # Load template if specified
            # Handle "null", "None", "undefined" strings as None
            if template_name and template_name not in ['null', 'None', 'undefined', '']:
                # User specified template
                template_filename = f"{template_name}.pptx"
                try:
                    template_bytes = ppt_manager.load_from_s3(template_filename)
                    ppt_manager.upload_to_code_interpreter(code_interpreter, template_filename, template_bytes)
                    logger.info(f"Using user template: {template_filename}")
                except FileNotFoundError:
                    code_interpreter.stop()
                    return {
                        "content": [{
                            "text": f"**Template not found**: {template_filename}"
                        }],
                        "status": "error"
                    }
            else:
                template_filename = None
                logger.info("Creating presentation from scratch")

            # Upload ppt_helpers (for utility functions)
            _upload_ppt_helpers_to_ci(code_interpreter)

            # Upload workspace images
            ppt_manager.load_workspace_images_to_ci(code_interpreter)

            # Generate creation code
            import json

            if slides:
                # Save slides to JSON file for safe transfer to Code Interpreter
                slides_json = json.dumps(slides, ensure_ascii=False, indent=2)
                slides_filename = f"slides_{presentation_name}.json"

                # Upload slides JSON file to Code Interpreter
                slides_upload_code = f"""
import json

# Save slides to file
slides_data = {repr(slides_json)}
with open('{slides_filename}', 'w', encoding='utf-8') as f:
    f.write(slides_data)

print(f"Slides file saved: {slides_filename}")
"""
                response = code_interpreter.invoke("executeCode", {
                    "code": slides_upload_code,
                    "language": "python",
                    "clearContext": False
                })

                # Check for errors
                for event in response.get("stream", []):
                    result = event.get("result", {})
                    if result.get("isError", False):
                        error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                        logger.error(f"Failed to upload slides: {error_msg[:200]}")
                        code_interpreter.stop()
                        return {
                            "content": [{"text": f"**Failed to prepare slides**: {error_msg[:500]}"}],
                            "status": "error"
                        }

                # Create presentation with custom_code for each slide
                code = f"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData
import json

# Load template or create blank
prs = Presentation({f"'{template_filename}'" if template_filename else ""})

# Set 16:9 aspect ratio (only for new presentations without template)
if not {f"'{template_filename}'" if template_filename else "None"}:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

slide_width = prs.slide_width
slide_height = prs.slide_height

# Load slides from file
with open('{slides_filename}', 'r', encoding='utf-8') as f:
    slides_data = json.load(f)

# Create each slide with custom_code
for i, slide_def in enumerate(slides_data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    custom_code = slide_def.get('custom_code', '')
    if custom_code:
        exec(custom_code)
        print(f"Slide {{i+1}}: created")
    else:
        print(f"Slide {{i+1}}: empty (no custom_code)")

# Save
prs.save('{presentation_filename}')

print(f"\\nCreated presentation with {{len(prs.slides)}} slides")
""".strip()
            else:
                # Create blank presentation with single slide
                code = f"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Create blank presentation
prs = Presentation({f"'{template_filename}'" if template_filename else ""})

# Set 16:9 aspect ratio (only for new presentations without template)
if not {f"'{template_filename}'" if template_filename else "None"}:
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

# Add blank slide if completely empty
if len(prs.slides) == 0:
    blank_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(blank_layout)

# Save
prs.save('{presentation_filename}')

print(f"Created blank presentation with {{len(prs.slides)}} slide(s)")
""".strip()

            response = code_interpreter.invoke("executeCode", {"code": code, "language": "python", "clearContext": False})

            # Check for errors
            for event in response.get("stream", []):
                result = event.get("result", {})
                if result.get("isError", False):
                    error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                    logger.error(f"Create presentation failed: {error_msg[:500]}")
                    code_interpreter.stop()
                    return {"content": [{"text": f"**Failed to create presentation**\n\n```\n{error_msg[:1000]}\n```"}], "status": "error"}

            # Download result
            output_ci_path = ppt_manager.get_ci_path(presentation_filename)
            file_bytes = ppt_manager.download_from_code_interpreter(code_interpreter, output_ci_path)

            if not file_bytes:
                code_interpreter.stop()
                return {"content": [{"text": "**Failed to retrieve presentation**"}], "status": "error"}

            # Save to S3
            s3_info = ppt_manager.save_to_s3(presentation_filename, file_bytes)

            # Save as artifact for Canvas display
            _save_ppt_artifact(
                tool_context=tool_context,
                filename=presentation_filename,
                s3_url=s3_info['s3_url'],
                size_kb=s3_info['size_kb'],
                tool_name='create_presentation',
                user_id=user_id,
                session_id=session_id
            )

            # Get updated workspace list
            documents = ppt_manager.list_s3_documents()
            other_files_count = len([d for d in documents if d['filename'] != presentation_filename])

            # Success message
            if slides:
                slide_count = len(slides)
                success_msg = f"""**Presentation created!**

**Filename:** {presentation_filename}
**Slides:** {slide_count}
**Template:** {template_filename or 'None (custom design)'}
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}

**Next steps:**
- Use `preview_presentation_slides` to check appearance
- Use `update_slide_content` to refine content
"""
            else:
                success_msg = f"""**Blank presentation created!**

**Filename:** {presentation_filename}
**Size:** {s3_info['size_kb']}
**Other files in workspace:** {other_files_count} presentation{'s' if other_files_count != 1 else ''}

**Next steps:**
- Use `add_slide` to add more slides with custom_code
"""

            code_interpreter.stop()
            return {
                "content": [{"text": success_msg}],
                "status": "success",
                "metadata": {
                    "filename": presentation_filename,
                    "slide_count": len(slides) if slides else 1,
                    "template": template_filename,
                    "size_kb": s3_info['size_kb'],
                    "s3_key": s3_info['s3_key'],
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id
                }
            }

        except Exception as e:
            code_interpreter.stop()
            raise e

    except Exception as e:
        logger.error(f"create_presentation error: {e}", exc_info=True)
        return {"content": [{"text": f"**Error:** {str(e)}"}], "status": "error"}


@tool(context=True)
def preview_presentation_slides(
    presentation_name: str,
    slide_numbers: list[int],
    tool_context: ToolContext
) -> Dict[str, Any]:
    """Get slide screenshots for YOU (the agent) to visually inspect before editing.

    This tool is for YOUR internal use - to see the actual layout, formatting,
    and content of slides before making modifications. Images are sent to you,
    not displayed to the user.

    Args:
        presentation_name: Presentation name without extension (e.g., "sales-deck")
        slide_numbers: List of slide numbers to preview (1-indexed).
                      Use empty list [] for all slides.
                      Example: [1, 3, 5] or []

    Use BEFORE modifying a presentation to:
    - See exact slide layout and formatting
    - Identify images, charts, or animations
    - Understand placeholder positions
    - Plan precise edits based on visual layout
    """
    import subprocess
    import tempfile
    import io
    from pdf2image import convert_from_path

    # Get user and session IDs
    user_id, session_id = _get_user_session_ids(tool_context)

    # Validate and prepare filename
    presentation_filename = f"{presentation_name}.pptx"
    logger.info(f"preview_presentation_slides: {presentation_filename}, slides {slide_numbers}")

    try:
        # Initialize presentation manager
        ppt_manager = PowerPointManager(user_id, session_id)

        # Check if presentation exists
        documents = ppt_manager.list_s3_documents()
        doc_info = next((d for d in documents if d['filename'] == presentation_filename), None)

        if not doc_info:
            available = [d['filename'] for d in documents if d['filename'].endswith('.pptx')]
            return {
                "content": [{
                    "text": f"Presentation not found: {presentation_filename}\n\n"
                           f"Available presentations: {', '.join(available) if available else 'None'}"
                }],
                "status": "error"
            }

        # Download presentation from S3
        pptx_bytes = ppt_manager.load_from_s3(presentation_filename)

        with tempfile.TemporaryDirectory() as temp_dir:
            # Save presentation to temp file
            pptx_path = os.path.join(temp_dir, presentation_filename)
            with open(pptx_path, 'wb') as f:
                f.write(pptx_bytes)

            # Convert PPTX to PDF using LibreOffice
            logger.info(f"Converting {presentation_filename} to PDF...")
            result = subprocess.run(
                ['soffice', '--headless', '--convert-to', 'pdf', '--outdir', temp_dir, pptx_path],
                capture_output=True,
                text=True,
                timeout=120  # 120 second timeout for large presentations
            )

            if result.returncode != 0:
                logger.error(f"LibreOffice conversion failed: {result.stderr}")
                return {
                    "content": [{
                        "text": f"PDF conversion failed\n\n{result.stderr}"
                    }],
                    "status": "error"
                }

            pdf_path = os.path.join(temp_dir, presentation_filename.replace('.pptx', '.pdf'))

            if not os.path.exists(pdf_path):
                return {
                    "content": [{
                        "text": "PDF file not created\n\nLibreOffice conversion may have failed silently."
                    }],
                    "status": "error"
                }

            # Get total pages in PDF (each slide becomes a page)
            from pdf2image import pdfinfo_from_path
            pdf_info = pdfinfo_from_path(pdf_path)
            total_slides = pdf_info.get('Pages', 1)

            # Determine which slides to preview
            if not slide_numbers:
                # Empty list means all slides
                target_slides = list(range(1, total_slides + 1))
            else:
                # Validate requested slide numbers
                invalid_slides = [s for s in slide_numbers if s < 1 or s > total_slides]
                if invalid_slides:
                    return {
                        "content": [{
                            "text": f"Invalid slide number(s): {invalid_slides}\n\n"
                                   f"Presentation has {total_slides} slides (1 to {total_slides})"
                        }],
                        "status": "error"
                    }
                target_slides = slide_numbers

            # Build content with images
            content = [{
                "text": f"**{presentation_filename}** - {len(target_slides)} slide(s) of {total_slides} total"
            }]

            for slide_num in target_slides:
                logger.info(f"Converting slide {slide_num} to image...")
                images = convert_from_path(
                    pdf_path,
                    first_page=slide_num,
                    last_page=slide_num,
                    dpi=150
                )

                if images:
                    img_buffer = io.BytesIO()
                    images[0].save(img_buffer, format='PNG')
                    img_bytes = img_buffer.getvalue()

                    content.append({"text": f"**Slide {slide_num}**"})
                    content.append({
                        "image": {
                            "format": "png",
                            "source": {"bytes": img_bytes}
                        }
                    })

            logger.info(f"Successfully generated {len(target_slides)} preview(s)")

            return {
                "content": content,
                "status": "success",
                "metadata": {
                    "filename": presentation_filename,
                    "slide_numbers": target_slides,
                    "total_slides": total_slides,
                    "tool_type": "powerpoint_presentation",
                    "user_id": user_id,
                    "session_id": session_id,
                    "hideImageInChat": True
                }
            }

    except subprocess.TimeoutExpired:
        logger.error("LibreOffice conversion timed out")
        return {
            "content": [{
                "text": "Conversion timed out\n\nThe presentation may be too large or complex."
            }],
            "status": "error"
        }
    except Exception as e:
        logger.error(f"preview_presentation_slides failed: {e}")
        return {
            "content": [{
                "text": f"Failed to generate preview\n\n{str(e)}"
            }],
            "status": "error"
        }


# --- Skill registration ---
register_skill("powerpoint-presentations", tools=[
    get_slide_code_examples, list_my_powerpoint_presentations, get_presentation_layouts,
    analyze_presentation, create_presentation, update_slide_content,
    add_slide, delete_slides, move_slide, duplicate_slide,
    update_slide_notes, preview_presentation_slides,
])
