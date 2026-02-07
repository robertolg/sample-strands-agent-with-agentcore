#!/usr/bin/env python3
"""
PowerPoint Slide Examples Test

Tests the slide_examples module and create_presentation hybrid format.

Usage:
    python scripts/test_ppt_examples.py
    python scripts/test_ppt_examples.py --execute    # Actually create presentation (uses API credits)
"""

import argparse
import sys
import os

# Add project source to path
sys.path.insert(0, os.path.join(os.path.dirname(__file__), '..', 'chatbot-app', 'agentcore', 'src'))

import boto3

# Configuration from environment
REGION = os.environ.get('AWS_REGION', 'us-west-2')
PROJECT_NAME = os.environ.get('PROJECT_NAME', 'strands-agent-chatbot')
ENVIRONMENT = os.environ.get('ENVIRONMENT', 'dev')


def get_code_interpreter_id() -> str:
    """Get Code Interpreter ID from environment or Parameter Store."""
    code_interpreter_id = os.getenv('CODE_INTERPRETER_ID')
    if code_interpreter_id:
        return code_interpreter_id

    try:
        ssm = boto3.client('ssm', region_name=REGION)
        param_name = f"/{PROJECT_NAME}/{ENVIRONMENT}/agentcore/code-interpreter-id"
        response = ssm.get_parameter(Name=param_name)
        return response['Parameter']['Value']
    except Exception as e:
        print(f"   Failed to get from SSM: {e}")
        return None


def test_slide_examples_import():
    """Test slide_examples module import."""
    print("\n📋 Test: Slide Examples Import")
    print("─" * 50)

    try:
        from builtin_tools.lib.slide_examples import get_examples, get_all_categories, SLIDE_EXAMPLES

        categories = get_all_categories()
        print(f"✅ slide_examples imported successfully")
        print(f"   Categories: {categories}")
        print(f"   Total categories: {len(categories)}")

        # Check each category has examples
        for cat in categories:
            examples = SLIDE_EXAMPLES[cat]['examples']
            print(f"   - {cat}: {len(examples)} example(s)")

        return True

    except ImportError as e:
        print(f"❌ Import failed: {e}")
        return False
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False


def test_ppt_utils_import():
    """Test ppt_utils module import."""
    print("\n🔧 Test: PPT Utils Import")
    print("─" * 50)

    try:
        from builtin_tools.lib.ppt_utils import (
            validate_presentation_name,
            sanitize_presentation_name,
            make_error_response,
            make_success_response,
        )

        # Test validation
        valid_result = validate_presentation_name("test-deck")
        invalid_result = validate_presentation_name("test_deck")

        print(f"✅ ppt_utils imported successfully")
        print(f"   validate_presentation_name('test-deck'): {valid_result}")
        print(f"   validate_presentation_name('test_deck'): {invalid_result}")

        return True

    except ImportError as e:
        print(f"❌ Import failed: {e}")
        return False
    except Exception as e:
        print(f"❌ Error: {e}")
        return False


def test_example_code_syntax():
    """Test that all example code has valid Python syntax."""
    print("\n🐍 Test: Example Code Syntax")
    print("─" * 50)

    try:
        from builtin_tools.lib.slide_examples import SLIDE_EXAMPLES

        all_valid = True
        for cat_name, cat_data in SLIDE_EXAMPLES.items():
            for example in cat_data['examples']:
                code = example['code']
                try:
                    compile(code, f"<{cat_name}/{example['name']}>", 'exec')
                    print(f"   ✅ {cat_name}/{example['name']}: Valid syntax")
                except SyntaxError as e:
                    print(f"   ❌ {cat_name}/{example['name']}: {e}")
                    all_valid = False

        if all_valid:
            print(f"✅ All example code has valid syntax")
        return all_valid

    except Exception as e:
        print(f"❌ Error: {e}")
        return False


def test_custom_code_execution(code_interpreter_id: str):
    """Test custom_code format execution in Code Interpreter."""
    print("\n🎯 Test: Custom Code Execution")
    print("─" * 50)

    try:
        from bedrock_agentcore.tools.code_interpreter_client import CodeInterpreter
        from builtin_tools.lib.slide_examples import SLIDE_EXAMPLES

        code_interpreter = CodeInterpreter(REGION)
        print(f"   Starting Code Interpreter...")
        code_interpreter.start(identifier=code_interpreter_id)

        # Get example code from text_layout category
        example_code = SLIDE_EXAMPLES['text_layout']['examples'][0]['code']

        # Wrap in presentation context
        test_code = f"""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

# Execute example code
{example_code}

prs.save('test_custom.pptx')
print("Custom code slide created successfully!")
"""

        print(f"   Executing custom code from text_layout example...")
        response = code_interpreter.invoke("executeCode", {
            "code": test_code,
            "language": "python",
            "clearContext": False
        })

        execution_success = False
        for event in response.get("stream", []):
            result = event.get("result", {})
            if result.get("isError", False):
                error_msg = result.get("structuredContent", {}).get("stderr", "Unknown error")
                print(f"❌ Execution failed: {error_msg[:500]}")
                code_interpreter.stop()
                return False

            stdout = result.get("structuredContent", {}).get("stdout", "")
            if stdout:
                print(f"   Output: {stdout}")
            execution_success = True

        if execution_success:
            print(f"✅ Custom code executed successfully!")

        code_interpreter.stop()
        return execution_success

    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False


def main():
    parser = argparse.ArgumentParser(description="Test PowerPoint Slide Examples")
    parser.add_argument("--execute", action="store_true", help="Actually execute code (uses API credits)")
    args = parser.parse_args()

    print("╔═══════════════════════════════════════════════════╗")
    print("║      PowerPoint Slide Examples Test               ║")
    print("╚═══════════════════════════════════════════════════╝")
    print()

    print(f"📍 Region: {REGION}")
    print(f"📁 Project: {PROJECT_NAME}")
    print(f"🌍 Environment: {ENVIRONMENT}")

    results = []

    # Test 1: slide_examples import
    results.append(("Slide Examples Import", test_slide_examples_import()))

    # Test 2: ppt_utils import
    results.append(("PPT Utils Import", test_ppt_utils_import()))

    # Test 3: Example code syntax
    results.append(("Example Code Syntax", test_example_code_syntax()))

    # Test 4 & 5: Execution tests (optional)
    if args.execute:
        code_interpreter_id = get_code_interpreter_id()
        if code_interpreter_id:
            print(f"\n⚠️  Running execution test (will use API credits)")
            print(f"   Code Interpreter ID: {code_interpreter_id}")
            results.append(("Custom Code Execution", test_custom_code_execution(code_interpreter_id)))
        else:
            print(f"\n❌ Code Interpreter ID not found, skipping execution tests")
    else:
        print("\n⏭️  Skipping execution tests (use --execute to enable)")

    # Summary
    print()
    print("═" * 50)
    print("📊 Test Summary")
    print("─" * 50)

    all_passed = True
    for name, passed in results:
        status = "✅" if passed else "❌"
        print(f"   {status} {name}")
        if not passed:
            all_passed = False

    print()
    if all_passed:
        print("✅ All tests passed!")
    else:
        print("⚠️  Some tests failed")
        sys.exit(1)


if __name__ == "__main__":
    main()
